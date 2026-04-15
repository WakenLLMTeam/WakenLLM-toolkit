#!/usr/bin/env python3
"""
analyst_slide.py — Single-slide, multi-panel analyst report generator.

Generates a dense, professional data-analysis PPTX slide with:
  • Header band  — title, subtitle, optional tag badge
  • Panel area   — 2–4 viz panels freely chosen by the LLM
  • Findings strip — 3–5 key takeaways at the bottom

Unlike slides_agent.py (which generates a full multi-slide deck), this script
is focused entirely on maximum information density within a SINGLE slide —
the kind of "executive one-pager" seen in McKinsey / BCG consulting deliverables
or sell-side equity research.

Layouts
-------
  left_main   55% left tall panel + 45% right with 2 stacked smaller panels  (3 panels)
  triptych    3 equal-width columns                                            (2 or 3 panels)
  dashboard   2×2 grid                                                         (4 panels)
  editorial   full-width top + 2 equal-width bottom panels                    (3 panels)

Usage
-----
  python analyst_slide.py \\
      --topic "Tesla FSD Chip Architecture & Performance" \\
      --output /tmp/tesla_fsd.pptx \\
      --theme bosch \\
      --panels 3

  python analyst_slide.py \\
      --topic "Bosch vs Mobileye radar sensor competitive analysis" \\
      --output /tmp/radar_analysis.pptx \\
      --panels 4 \\
      --lang zh

Environment variables (LLM backend — same as llm_planner.py):
  ANTHROPIC_API_KEY / ANTHROPIC_AUTH_TOKEN  → Claude
  OPENAI_API_KEY                            → OpenAI
  CUSTOM_LLM_BASE_URL + CUSTOM_LLM_API_KEY  → custom OpenAI-compatible
  OLLAMA_HOST / OLLAMA_MODEL                → Ollama local
"""
from __future__ import annotations

import argparse
import json
import os
import sys
import textwrap
import traceback
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# Allow importing sibling scripts without installing as a package
_SCRIPTS_DIR = Path(__file__).parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

import math as _math

import matplotlib
matplotlib.use("Agg")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

import llm_planner as _lp
import build_deck as _bd
import build_pptx as _bpx
import slides_agent as _sa
import viz_theme


# ── Slide geometry constants ──────────────────────────────────────────────────

_SLIDE_W = 13.333   # inches  (16:9)
_SLIDE_H = 7.5      # inches

_HEADER_H    = 0.82  # header band height (title + subtitle)
_FINDINGS_H  = 1.05  # findings strip height at bottom
_PANEL_TOP   = _HEADER_H + 0.10   # top of panel area
_PANEL_BOT   = _SLIDE_H - _FINDINGS_H - 0.08
_PANEL_H     = _PANEL_BOT - _PANEL_TOP  # available panel height ≈ 5.55"
_MARGIN      = 0.18  # left/right margin
_GAP         = 0.12  # gap between panels

_PANEL_AREA_W = _SLIDE_W - 2 * _MARGIN  # ≈ 12.97"


# ── Theme ─────────────────────────────────────────────────────────────────────

_THEMES: Dict[str, Dict[str, Any]] = {
    "bosch": {
        "accent_rgb":   [226,  0, 21],
        "header_bg":    [226,  0, 21],
        "header_text":  [255, 255, 255],
        "body_rgb":     [ 17, 17, 17],
        "findings_bg":  [ 34, 34, 34],
        "findings_text":[255, 255, 255],
        "panel_border": [200, 211, 224],
        "tag_bg":       [255, 255, 255],
        "tag_text":     [226,  0, 21],
        "font":         "Calibri",
    },
    "default": {
        "accent_rgb":   [ 37, 99, 235],
        "header_bg":    [ 37, 99, 235],
        "header_text":  [255, 255, 255],
        "body_rgb":     [ 17, 17, 17],
        "findings_bg":  [ 30, 41, 59],
        "findings_text":[255, 255, 255],
        "panel_border": [200, 211, 224],
        "tag_bg":       [255, 255, 255],
        "tag_text":     [ 37, 99, 235],
        "font":         "Calibri",
    },
}


def _tc(cfg: Dict, key: str) -> RGBColor:
    v = cfg[key]
    return RGBColor(v[0], v[1], v[2])


# ── Layout engine ─────────────────────────────────────────────────────────────
#
# Layout catalogue
# ----------------
# 2-panel:
#   duo          — 2 equal side-by-side columns
#
# 3-panel:
#   hero_left    — panel 0 tall left (55%), panels 1+2 stacked right (45%)
#   hero_right   — panels 0+1 stacked left (45%), panel 2 tall right (55%)
#   hero_top     — panel 0 full-width top (~50%), panels 1+2 equal bottom
#   hero_bottom  — panels 0+1 equal top, panel 2 full-width bottom (~50%)
#   triptych     — 3 equal columns
#
# 4-panel:
#   dashboard    — 2×2 equal grid
#   wide_top     — panel 0 full-width top (45%), panels 1+2+3 equal columns below
#
# Legacy aliases kept for backward compatibility:
#   left_main  → hero_left
#   editorial  → hero_top

_ALL_LAYOUTS = {
    "duo", "hero_left", "hero_right", "hero_top", "hero_bottom",
    "triptych", "dashboard", "wide_top",
    # aliases
    "left_main", "editorial",
}

# Wide viz types that benefit from a full-width hero panel
_WIDE_VIZ = {"timeline", "gantt", "pipeline", "comparison", "arch", "bar_chart",
             "line_chart", "waterfall"}
# Roughly-square viz types that work well in smaller grid cells
_SQUARE_VIZ = {"radar", "matrix_2x2", "venn", "onion", "pie", "funnel",
               "scatter", "heatmap", "tree", "swot"}


def _auto_layout(panels: List[Dict[str, Any]], n_panels: int) -> str:
    """
    Infer the best layout when the LLM did not return a valid one.

    Rules (checked in priority order):
    1. hero-importance panel present + n==3 → hero_left or hero_top
       (hero_top if the hero viz type is "wide"; hero_left otherwise)
    2. Wide viz in panel 0 + n==3 → hero_top
    3. Wide viz in panel 0 + n==4 → wide_top
    4. n==2 → duo
    5. n==3 → triptych
    6. n==4 → dashboard
    """
    hero_idx = next(
        (p.get("panel_index", i)
         for i, p in enumerate(panels)
         if p.get("importance") == "hero"),
        None,
    )
    def _vtype(idx):
        for p in panels:
            if p.get("panel_index", 0) == idx:
                return (p.get("viz") or {}).get("type", "")
        return ""

    if n_panels == 2:
        return "duo"
    if n_panels == 3:
        if hero_idx == 0 and _vtype(0) in _WIDE_VIZ:
            return "hero_top"
        if hero_idx == 0:
            return "hero_left"
        if hero_idx == 2:
            return "hero_right"
        if _vtype(0) in _WIDE_VIZ:
            return "hero_top"
        return "triptych"
    if n_panels == 4:
        if hero_idx == 0 and _vtype(0) in _WIDE_VIZ:
            return "wide_top"
        return "dashboard"
    return "triptych"


def _build_panel_boxes(layout: str, n_panels: int) -> List[Tuple[float, float, float, float]]:
    """
    Return list of (left_in, top_in, w_in, h_in) for each panel.

    Panels occupy:
      x: [_MARGIN, _SLIDE_W - _MARGIN]   (width = _PANEL_AREA_W)
      y: [_PANEL_TOP, _PANEL_BOT]        (height = _PANEL_H)
    """
    # Resolve legacy aliases
    if layout == "left_main":
        layout = "hero_left"
    elif layout == "editorial":
        layout = "hero_top"

    L = _MARGIN
    T = _PANEL_TOP
    W = _PANEL_AREA_W
    H = _PANEL_H
    G = _GAP

    # ── 2-panel ──────────────────────────────────────────────────────────────
    if layout == "duo":
        col_w = (W - G) / 2
        return [
            (L,              T, col_w, H),   # panel 0 — left
            (L + col_w + G,  T, col_w, H),   # panel 1 — right
        ][:n_panels]

    # ── 3-panel hero variants ─────────────────────────────────────────────────
    elif layout == "hero_left":
        lw = W * 0.55 - G / 2
        rw = W * 0.45 - G / 2
        rh = (H - G) / 2
        return [
            (L,          T,              lw, H),    # 0 — tall left hero
            (L + lw + G, T,              rw, rh),   # 1 — top-right detail
            (L + lw + G, T + rh + G,    rw, rh),   # 2 — bot-right detail
        ][:n_panels]

    elif layout == "hero_right":
        lw = W * 0.45 - G / 2
        rw = W * 0.55 - G / 2
        lh = (H - G) / 2
        return [
            (L,          T,              lw, lh),   # 0 — top-left detail
            (L,          T + lh + G,    lw, lh),   # 1 — bot-left detail
            (L + lw + G, T,              rw, H),    # 2 — tall right hero
        ][:n_panels]

    elif layout == "hero_top":
        th = H * 0.50
        bh = H - th - G
        cw = (W - G) / 2
        return [
            (L,          T,           W,  th),   # 0 — full-width top hero
            (L,          T + th + G, cw, bh),   # 1 — bot-left detail
            (L + cw + G, T + th + G, cw, bh),   # 2 — bot-right detail
        ][:n_panels]

    elif layout == "hero_bottom":
        th = H * 0.50
        bh = H - th - G
        cw = (W - G) / 2
        return [
            (L,          T,  cw, th),   # 0 — top-left detail
            (L + cw + G, T,  cw, th),   # 1 — top-right detail
            (L,          T + th + G, W, bh),    # 2 — full-width bottom hero
        ][:n_panels]

    elif layout == "triptych":
        cols = max(2, min(n_panels, 3))
        cw = (W - G * (cols - 1)) / cols
        return [
            (L + i * (cw + G), T, cw, H)
            for i in range(cols)
        ][:n_panels]

    # ── 4-panel ───────────────────────────────────────────────────────────────
    elif layout == "dashboard":
        cw = (W - G) / 2
        rh = (H - G) / 2
        return [
            (L + (i % 2) * (cw + G),
             T + (i // 2) * (rh + G),
             cw, rh)
            for i in range(min(n_panels, 4))
        ]

    elif layout == "wide_top":
        th = H * 0.45
        bh = H - th - G
        cols = 3
        cw = (W - G * (cols - 1)) / cols
        boxes = [(L, T, W, th)]   # panel 0 — full-width hero
        for i in range(cols):
            boxes.append((L + i * (cw + G), T + th + G, cw, bh))
        return boxes[:n_panels]

    else:
        # Unknown layout — auto-infer
        return _build_panel_boxes(_auto_layout([], n_panels), n_panels)


# ── Viz rendering ─────────────────────────────────────────────────────────────

# Minimum safe render sizes (width, height in inches) per chart type.
# Charts are rendered at these sizes then letterboxed into the panel box,
# rather than being forced into the panel's exact dimensions.
_MIN_VIZ_SIZE: Dict[str, Tuple[float, float]] = {
    "flowchart":  (8.0, 5.0),
    "arch":       (8.0, 4.5),
    "tree":       (7.5, 5.0),
    "timeline":   (9.0, 3.5),
    "matrix_2x2": (6.5, 5.5),
    "scatter":    (6.5, 5.5),
    "heatmap":    (7.0, 5.0),
    "comparison": (8.0, 4.5),
    "radar":      (6.0, 5.5),
    "venn":       (6.5, 5.5),
    "onion":      (6.0, 6.0),
    "gantt":      (10.0, 4.0),
    "swot":       (7.5, 5.5),
    "pipeline":   (9.0, 3.5),
    "bar_chart":  (7.0, 4.5),
    "line_chart": (7.0, 4.5),
    "waterfall":  (8.0, 4.5),
    "funnel":     (5.5, 5.5),
    "pie":        (5.5, 5.0),
}
_DEFAULT_MIN_VIZ_SIZE: Tuple[float, float] = (7.0, 4.5)


def _render_panels(
    plan: Dict[str, Any],
    assets_dir: str,
    panel_boxes: List[Tuple[float, float, float, float]],
) -> List[Optional[str]]:
    """
    Render each panel's viz spec to a PNG.  Returns a list of paths (or None on
    failure) in panel_index order.
    """
    Path(assets_dir).mkdir(parents=True, exist_ok=True)
    panels = plan.get("panels", [])
    png_paths: List[Optional[str]] = [None] * len(panels)

    for panel in panels:
        idx = int(panel.get("panel_index", 0))
        viz = panel.get("viz")
        if not isinstance(viz, dict):
            continue

        # Normalize & validate (reuse slides_agent helpers); up to 2 repair passes
        viz = _sa._normalize_viz_spec(viz)
        panel["viz"] = viz
        err = _sa._validate_viz_spec(viz)
        for _repair_pass in range(3):
            if not err:
                break
            print(f"[analyst_slide] panel {idx} viz invalid (pass {_repair_pass+1}): {err} — LLM repair…",
                  file=sys.stderr)
            repaired = _bd._repair_spec(viz, f"Validation error: {err}")
            if repaired is None:
                break
            repaired = _sa._normalize_viz_spec(repaired)
            err2 = _sa._validate_viz_spec(repaired)
            viz = repaired
            panel["viz"] = viz
            err = err2
            if not err:
                print(f"[analyst_slide] panel {idx} repaired OK (pass {_repair_pass+1})", file=sys.stderr)
        if err:
            print(f"[analyst_slide] panel {idx} giving up after repairs: {err}", file=sys.stderr)
            continue

        viz_type = viz.get("type", "")
        renderer  = _bd._RENDERERS.get(viz_type)
        if renderer is None:
            print(f"[analyst_slide] panel {idx}: no renderer for '{viz_type}'", file=sys.stderr)
            continue

        # Inject figure size using the chart's minimum safe dimensions.
        # We render at natural quality then _add_picture_fit letterboxes into panel.
        # Never force the chart into the panel's exact (potentially tiny) size.
        box = panel_boxes[idx] if idx < len(panel_boxes) else None
        if box:
            min_w, min_h = _MIN_VIZ_SIZE.get(viz_type, _DEFAULT_MIN_VIZ_SIZE)
            # Allow larger if the panel is bigger than the minimum
            render_w = max(min_w, box[2] * 0.95)
            render_h = max(min_h, box[3] * 0.95)
            viz["fig_width"]  = round(render_w, 2)
            viz["fig_height"] = round(render_h, 2)

        out_png = os.path.join(assets_dir, f"panel_{idx}_{viz_type}.png")

        # Use build_deck's retry-with-repair logic
        final_spec = _bd._render_with_retry(
            renderer, viz, out_png, idx, viz_type, max_retries=2,
        )
        if final_spec is not None:
            panel["viz"] = final_spec
            png_paths[idx] = out_png

    return png_paths


# ── PPTX assembly ─────────────────────────────────────────────────────────────

def _add_textbox(
    slide,
    text: str,
    left: float, top: float, w: float, h: float,
    *,
    font_name: str,
    font_pt: float,
    bold: bool = False,
    color: RGBColor,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    word_wrap: bool = True,
):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _assemble_slide(
    plan: Dict[str, Any],
    panel_pngs: List[Optional[str]],
    panel_boxes: List[Tuple[float, float, float, float]],
    output_path: str,
    theme_name: str = "bosch",
) -> str:
    """Build the single-slide PPTX and return the output path."""
    cfg = _THEMES.get(theme_name, _THEMES["default"])
    font = cfg["font"]

    prs = Presentation()
    prs.slide_width  = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)

    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Header band ──────────────────────────────────────────────────────────
    from pptx.util import Inches as _I
    from pptx.dml.color import RGBColor as _RGB
    import pptx.util as _pu

    hdr = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        _I(0), _I(0), _I(_SLIDE_W), _I(_HEADER_H),
    )
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = _tc(cfg, "header_bg")
    hdr.line.fill.background()

    title_text    = _bpx._truncate(plan.get("title", "Analysis"), 60)
    subtitle_text = _bpx._truncate(plan.get("subtitle", ""), 90)
    tag_text      = _bpx._truncate(plan.get("tag", ""), 18)

    title_pt = _bpx._fit_text_pt(title_text, _SLIDE_W - 3.0, _HEADER_H * 0.55, start_pt=26, min_pt=14)
    _add_textbox(
        slide, title_text,
        left=_MARGIN, top=0.03, w=_SLIDE_W - 3.0, h=_HEADER_H * 0.58,
        font_name=font, font_pt=title_pt, bold=True,
        color=_tc(cfg, "header_text"),
        align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE,
    )

    if subtitle_text:
        sub_pt = _bpx._fit_text_pt(subtitle_text, _SLIDE_W - 3.2, _HEADER_H * 0.40, start_pt=13, min_pt=9)
        _add_textbox(
            slide, subtitle_text,
            left=_MARGIN, top=_HEADER_H * 0.55, w=_SLIDE_W - 3.2, h=_HEADER_H * 0.42,
            font_name=font, font_pt=sub_pt, bold=False,
            color=_tc(cfg, "header_text"),
            align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP,
        )

    if tag_text:
        tag_w = min(len(tag_text) * 0.14 + 0.5, 2.0)
        tag_l = _SLIDE_W - _MARGIN - tag_w
        tag_shape = slide.shapes.add_shape(
            1,
            Inches(tag_l), Inches(0.18), Inches(tag_w), Inches(0.44),
        )
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = _tc(cfg, "tag_bg")
        tag_shape.line.color.rgb = _tc(cfg, "tag_text")
        tag_shape.line.width = Pt(1.2)
        tf = tag_shape.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag_text
        run.font.name = font
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = _tc(cfg, "tag_text")

    # ── Panel separator line ──────────────────────────────────────────────────
    sep = slide.shapes.add_shape(
        1,
        Inches(0), Inches(_HEADER_H),
        Inches(_SLIDE_W), Inches(0.03),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = _tc(cfg, "panel_border")
    sep.line.fill.background()

    # ── Panels ───────────────────────────────────────────────────────────────
    panels = plan.get("panels", [])
    for panel in panels:
        idx = int(panel.get("panel_index", 0))
        if idx >= len(panel_boxes):
            continue
        box_l, box_t, box_w, box_h = panel_boxes[idx]

        # Panel background card
        card = slide.shapes.add_shape(
            1,
            Inches(box_l), Inches(box_t), Inches(box_w), Inches(box_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 251, 252)
        card.line.color.rgb = _tc(cfg, "panel_border")
        card.line.width = Pt(0.6)

        # Caption strip at top of panel
        cap_h = 0.28
        caption = _bpx._truncate(panel.get("caption", ""), 50)
        if caption:
            cap_pt = _bpx._fit_text_pt(caption, box_w - 0.16, cap_h, start_pt=11, min_pt=8)
            _add_textbox(
                slide, caption,
                left=box_l + 0.08, top=box_t + 0.04, w=box_w - 0.16, h=cap_h,
                font_name=font, font_pt=cap_pt, bold=True,
                color=_tc(cfg, "body_rgb"),
                align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP,
            )

        # Image (if rendered); fallback to placeholder text on failure
        png = panel_pngs[idx] if idx < len(panel_pngs) else None
        img_top = box_t + cap_h + 0.06 if caption else box_t + 0.04
        img_h   = box_h - (img_top - box_t) - 0.04
        if png and Path(png).exists():
            try:
                _bpx._add_picture_fit(
                    slide, png,
                    Inches(box_l + 0.04), Inches(img_top),
                    Inches(box_w - 0.08), Inches(img_h),
                )
            except Exception as e:
                print(f"[analyst_slide] panel {idx} image insert failed: {e}", file=sys.stderr)
                png = None  # fall through to placeholder
        if not (png and Path(png).exists()):
            # Placeholder — show viz type so it's clear what was intended
            viz_lbl = (panel.get("viz") or {}).get("type", "chart")
            _add_textbox(
                slide, f"[ {viz_lbl} — rendering unavailable ]",
                left=box_l + 0.12, top=img_top + img_h * 0.35,
                w=box_w - 0.24, h=img_h * 0.3,
                font_name=font, font_pt=11, bold=False,
                color=RGBColor(160, 160, 160),
                align=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    # ── Findings strip ────────────────────────────────────────────────────────
    f_top = _SLIDE_H - _FINDINGS_H
    findings_bg = slide.shapes.add_shape(
        1,
        Inches(0), Inches(f_top),
        Inches(_SLIDE_W), Inches(_FINDINGS_H),
    )
    findings_bg.fill.solid()
    findings_bg.fill.fore_color.rgb = _tc(cfg, "findings_bg")
    findings_bg.line.fill.background()

    # "KEY FINDINGS" label
    _add_textbox(
        slide, "KEY FINDINGS",
        left=_MARGIN, top=f_top + 0.06, w=1.4, h=0.28,
        font_name=font, font_pt=9, bold=True,
        color=_tc(cfg, "findings_text"),
        align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.TOP,
    )

    findings = (plan.get("findings") or [])[:5]
    if findings:
        n = len(findings)
        avail_w = _SLIDE_W - _MARGIN - 1.6  # after label
        col_w   = avail_w / n
        for fi, finding in enumerate(findings):
            ft = _bpx._truncate(str(finding), 80)
            fpt = _bpx._fit_text_pt(ft, col_w - 0.2, _FINDINGS_H - 0.18, start_pt=11, min_pt=8)
            _add_textbox(
                slide, f"• {ft}",
                left=_MARGIN + 1.6 + fi * col_w,
                top=f_top + 0.06,
                w=col_w - 0.12,
                h=_FINDINGS_H - 0.14,
                font_name=font, font_pt=fpt, bold=False,
                color=_tc(cfg, "findings_text"),
                align=PP_ALIGN.LEFT, vertical_anchor=MSO_ANCHOR.MIDDLE,
            )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


# ── LLM planning ──────────────────────────────────────────────────────────────

_ANALYST_SYSTEM = textwrap.dedent("""
You are an elite data analyst and visual storytelling expert. Given a topic and
the number of analysis panels desired, you produce a single dense "analyst slide"
plan as a JSON object.

CRITICAL RULES
- Return ONLY valid JSON — no markdown fences, no commentary.
- title      : ≤52 characters, punchy headline
- subtitle   : ≤85 characters, key context or time range
- tag        : ≤14 characters, e.g. "CONFIDENTIAL", "DRAFT", "Q2 2025" (optional)
- layout     : choose from the layout catalogue below
- panels     : list of panel objects in panel_index order (0-based)
- findings   : 3–5 bullet takeaways, each ≤75 characters, action-oriented

━━━ LAYOUT CATALOGUE ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

2-panel layouts:
  "duo"           — [0|1] two equal side-by-side panels

3-panel layouts:
  "hero_left"     — [0(tall)|[1,2 stacked]]  panel 0 wide (55%), 1+2 small stacked right
  "hero_right"    — [[0,1 stacked]|2(tall)]  panels 0+1 small stacked left, panel 2 wide (55%)
  "hero_top"      — [0(full-width)] / [1|2]  panel 0 spans full width (top), 1+2 detail below
  "hero_bottom"   — [0|1] / [2(full-width)]  panels 0+1 detail above, panel 2 spans full width
  "triptych"      — [0|1|2] three equal columns

4-panel layouts:
  "dashboard"     — [[0|1],[2|3]]  2×2 equal grid
  "wide_top"      — [0(full-width)] / [1|2|3]  panel 0 spans full width, 1+2+3 three columns

━━━ ADAPTIVE LAYOUT SELECTION RULES ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Step 1 — identify panel importance:
  Assign "importance": "hero" to the most information-rich or visually dominant panel,
  and "importance": "detail" to supporting panels.

Step 2 — match to layout:

  IF n_panels == 2:
    → "duo"  (always)

  IF n_panels == 3:
    • hero panel 0 AND its viz type is WIDE  → "hero_top"
      (WIDE = timeline, gantt, pipeline, comparison, arch, bar_chart, line_chart, waterfall)
    • hero panel 0 AND viz type is tall/square → "hero_left"
    • hero panel 2 (end of sequence)          → "hero_right"
    • hero panel 2 AND viz type is WIDE        → "hero_bottom"
    • all panels equally important             → "triptych"

  IF n_panels == 4:
    • hero panel 0 AND viz type is WIDE  → "wide_top"
    • otherwise                          → "dashboard"

━━━ PANEL OBJECT SCHEMA ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{
  "panel_index": 0,
  "importance": "hero" | "detail",
  "caption": "≤45 chars — descriptive chart title",
  "viz": { ...full viz spec... }
}

SUPPORTED VIZ TYPES (choose the most analytically appropriate for each panel)
  bar_chart, line_chart, scatter, heatmap, waterfall, funnel, pie, radar,
  comparison, pipeline, timeline, arch, tree, matrix_2x2, venn, onion,
  gantt, swot, flowchart

MANDATORY COMPLETE SCHEMAS FOR COMPLEX TYPES (copy exactly, fill in your data):

comparison — ALL four keys REQUIRED (rows, cols, cells, highlight_col):
  {"type":"comparison","title":"...","rows":["row1","row2","row3"],"cols":["ColA","ColB","ColC"],"highlight_col":1,"cells":[["a1","b1","c1"],["a2","b2","c2"],["a3","b3","c3"]],"row_notes":["","",""]  }

heatmap — ALL four keys REQUIRED (rows, cols, values, show_values):
  {"type":"heatmap","title":"...","rows":["R1","R2","R3"],"cols":["C1","C2","C3"],"values":[[8,5,3],[6,9,4],[7,2,8]],"color_scheme":"blue","show_values":true}

scatter — REQUIRED: series with at least one point:
  {"type":"scatter","title":"...","x_label":"X →","y_label":"↑ Y","series":[{"name":"Group","points":[{"x":3,"y":7,"label":"A","size":200},{"x":7,"y":5,"label":"B","size":150}]}]}

timeline — REQUIRED: non-empty stages list:
  {"type":"timeline","title":"...","stages":[{"label":"Phase 1","year":"2020","annotation":"desc"},{"label":"Phase 2","year":"2022","annotation":"desc"}]}

matrix_2x2 — REQUIRED: non-empty items list with x,y in [0,1]:
  {"type":"matrix_2x2","title":"...","x_label":"X →","y_label":"↑ Y","quadrants":{"top_left":{"label":"A"},"top_right":{"label":"B"},"bottom_left":{"label":"C"},"bottom_right":{"label":"D"}},"items":[{"label":"Item1","x":0.3,"y":0.7},{"label":"Item2","x":0.8,"y":0.6}]}

arch — REQUIRED: layers with non-empty blocks in EVERY layer (no empty blocks arrays):
  {"type":"arch","title":"...","direction":"BT","layers":[{"name":"Application","color":"#dbeafe","blocks":[{"label":"Path Planner","sublabel":"Behavior · Route","badge":"ASIL-D"},{"label":"AEB","sublabel":"Emergency braking"}]},{"name":"Middleware","color":"#fef9c3","blocks":[{"label":"ROS2","sublabel":"DDS transport"},{"label":"AUTOSAR","sublabel":"COM stack"}]},{"name":"Hardware","color":"#f3e8ff","blocks":[{"label":"SoC","sublabel":"Orin / EyeQ6"},{"label":"MCU","sublabel":"Safety core"}]}]}
  CRITICAL: every layer object MUST have "blocks" as a non-empty list — never leave "blocks": []

flowchart — REQUIRED: nodes with id+label+shape+color, edges with from+to:
  {"type":"flowchart","title":"...","layout":"TB","nodes":[{"id":"n1","label":"Input","shape":"rect","color":"#dbeafe"},{"id":"n2","label":"Process","shape":"rect","color":"#dcfce7"},{"id":"n3","label":"Decision","shape":"diamond","color":"#fef9c3"},{"id":"n4","label":"Output","shape":"rounded","color":"#f3e8ff"}],"edges":[{"from":"n1","to":"n2"},{"from":"n2","to":"n3","label":"check"},{"from":"n3","to":"n4","label":"pass"}]}

KEY PRINCIPLE — INFORMATION DENSITY
Each panel should encode a distinct analytical dimension of the topic:
  • Panel 0 (dominant/wide): macro trend, competitive benchmark, or architecture
  • Panel 1: breakdown, composition, or process flow
  • Panel 2: comparison, heatmap, or risk assessment
  • Panel 3 (if dashboard): forward-looking or actionable view

Do NOT repeat the same viz type across panels unless the data is genuinely different.
Prefer quantified data over qualitative lists. Include realistic plausible numbers.

OUTPUT SCHEMA
{
  "title": "string",
  "subtitle": "string",
  "tag": "string or null",
  "layout": "duo|hero_left|hero_right|hero_top|hero_bottom|triptych|dashboard|wide_top",
  "panels": [
    {
      "panel_index": 0,
      "importance": "hero|detail",
      "caption": "string",
      "viz": { "type": "...", ...full spec... }
    }
  ],
  "findings": ["string", ...]
}
""").strip()

_ANALYST_USER = textwrap.dedent("""
Topic: {topic}
Number of panels: {n_panels}
Language: {lang_instruction}
Theme: {theme}

Generate the analyst slide plan JSON now. Make the data realistic, specific, and
analytically compelling for an audience of senior executives or domain experts.
""").strip()

_LANG_INSTRUCTIONS = {
    "zh": "Chinese (Simplified) — all labels, captions, findings, title, subtitle in Chinese",
    "en": "English — all text in English",
}


def _call_analyst_llm(
    topic: str,
    n_panels: int,
    theme: str,
    lang: str,
    backend: str,
    model: str,
) -> Dict[str, Any]:
    lang_instruction = _LANG_INSTRUCTIONS.get(lang, _LANG_INSTRUCTIONS["en"])
    user = _ANALYST_USER.format(
        topic=topic,
        n_panels=n_panels,
        lang_instruction=lang_instruction,
        theme=theme,
    )
    raw = _lp._call_llm(_ANALYST_SYSTEM, user, backend, model)
    return _lp._extract_json(raw)


# ── Main pipeline ─────────────────────────────────────────────────────────────

def build_analyst_slide(
    topic: str,
    output_path: str,
    *,
    theme: str = "bosch",
    lang: str = "en",
    n_panels: int = 3,
    backend: str = "",
    model: str = "",
    assets_dir: Optional[str] = None,
    plan: Optional[Dict[str, Any]] = None,
) -> str:
    """
    Full pipeline: topic → LLM plan → render panels → assemble PPTX.
    Returns the output path.
    """
    if not backend:
        backend = _lp._detect_backend()
    if not model:
        model = os.environ.get("CUSTOM_LLM_MODEL", "") if backend == "custom" else ""

    n_panels = max(2, min(4, int(n_panels)))

    # Stage 1: get plan from LLM (or use supplied plan)
    if plan is None:
        print(f"[analyst_slide] Calling LLM ({backend}) for analyst plan…")
        plan = _call_analyst_llm(topic, n_panels, theme, lang, backend, model)
        print(f"[analyst_slide] Plan: layout={plan.get('layout')}, "
              f"panels={len(plan.get('panels', []))}, "
              f"findings={len(plan.get('findings', []))}")
    else:
        print("[analyst_slide] Using supplied plan.")

    actual_n = len(plan.get("panels", []))
    if actual_n == 0:
        raise ValueError("LLM returned no panels in plan")

    # Choose layout from plan (or auto-infer from panel importance + viz types)
    layout = plan.get("layout", "")
    if layout not in _ALL_LAYOUTS:
        layout = _auto_layout(plan.get("panels", []), actual_n)
        plan["layout"] = layout
        print(f"[analyst_slide] Layout auto-inferred: {layout}")
    else:
        # Resolve legacy aliases
        if layout == "left_main":
            layout = "hero_left"
        elif layout == "editorial":
            layout = "hero_top"
        plan["layout"] = layout

    panel_boxes = _build_panel_boxes(layout, actual_n)

    # Stage 2: render panels
    if assets_dir is None:
        stem = Path(output_path).stem
        assets_dir = str(Path(output_path).parent / f"{stem}_assets")

    print(f"[analyst_slide] Rendering {actual_n} panels ({layout} layout)…")
    panel_pngs = _render_panels(plan, assets_dir, panel_boxes)
    ok = sum(1 for p in panel_pngs if p is not None)
    print(f"[analyst_slide] {ok}/{actual_n} panels rendered successfully.")

    # Stage 3: assemble PPTX
    print(f"[analyst_slide] Assembling PPTX → {output_path}")
    _assemble_slide(plan, panel_pngs, panel_boxes, output_path, theme_name=theme)
    print(f"[analyst_slide] Done: {output_path}")
    return output_path


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> int:
    parser = argparse.ArgumentParser(
        description="Generate a single dense analyst-report slide as PPTX"
    )
    parser.add_argument("--topic",   required=True, help="Analysis topic")
    parser.add_argument("--output",  required=True, help="Output .pptx path")
    parser.add_argument("--theme",   default="bosch", choices=list(_THEMES),
                        help="Visual theme (default: bosch)")
    parser.add_argument("--lang",    default="en", choices=["en", "zh"],
                        help="Output language (default: en)")
    parser.add_argument("--panels",  type=int, default=3, choices=[2, 3, 4],
                        help="Number of viz panels (default: 3)")
    parser.add_argument("--backend", default="",
                        choices=["", "anthropic", "openai", "custom", "ollama"],
                        help="LLM backend (auto-detected if omitted)")
    parser.add_argument("--model",   default="",
                        help="Model name (backend-specific; uses env default if omitted)")
    parser.add_argument("--assets-dir", default=None,
                        help="Directory to store panel PNG assets")
    parser.add_argument("--plan-file", default=None,
                        help="Skip LLM call and use an existing plan JSON file")
    parser.add_argument("--plan-only", action="store_true",
                        help="Print the LLM plan JSON and exit without building PPTX")
    args = parser.parse_args()

    # Load pre-made plan if supplied
    plan = None
    if args.plan_file:
        with open(args.plan_file, encoding="utf-8") as f:
            plan = json.load(f)

    if args.plan_only:
        if plan is None:
            backend = args.backend or _lp._detect_backend()
            model   = args.model or (os.environ.get("CUSTOM_LLM_MODEL", "") if backend == "custom" else "")
            plan = _call_analyst_llm(args.topic, args.panels, args.theme, args.lang, backend, model)
        print(json.dumps(plan, ensure_ascii=False, indent=2))
        return 0

    build_analyst_slide(
        args.topic,
        args.output,
        theme=args.theme,
        lang=args.lang,
        n_panels=args.panels,
        backend=args.backend,
        model=args.model,
        assets_dir=args.assets_dir,
        plan=plan,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
