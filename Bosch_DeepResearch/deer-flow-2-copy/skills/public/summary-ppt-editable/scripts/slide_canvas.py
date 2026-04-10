#!/usr/bin/env python3
"""
slide_canvas.py — LEGO-style slide builder.

Agent composes a slide by placing components (text, viz, image, shape)
onto a Canvas, then calls finish() to render to PPTX.

Usage (API):
    from slide_canvas import Canvas
    cv = Canvas(width_in=13.333, height_in=7.5)
    cv.add_text(id="t1", content="Title", pos=(0.5, 0.4), size=(8, 0.7),
                font_size=28, bold=True)
    cv.add_viz(id="v1", viz_type="timeline",
               viz_data={"stages": [{"label": "L2", "year": "2018"}]},
               pos=(0.5, 1.3), size=(12, 4.5))
    cv.finish("/tmp/slide1.pptx")
"""
from __future__ import annotations

import copy
import json
import os
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Allow importing sibling renderers / theme without installing as a package
_SCRIPTS_DIR = Path(__file__).parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from viz_theme import THEME, BOSCH_THEME, setup_matplotlib
from build_deck import _RENDERERS

setup_matplotlib()


# ── Dataclasses ──────────────────────────────────────────────────────────────

@dataclass
class TextComponent:
    id: str
    content: str
    pos: tuple[float, float]          # (x, y) inches from top-left
    size: tuple[float, float]           # (w, h) inches
    font_size: float = 12.0
    bold: bool = False
    color_rgb: Optional[list[int]] = None
    alignment: str = "left"            # left | center | right
    bullet: bool = False
    level: int = 0                     # 0 = top-level bullet, 1 = sub-bullet
    font_name: Optional[str] = None
    wrap: bool = True
    italic: bool = False
    line_spacing: Optional[float] = None


@dataclass
class ShapeComponent:
    id: str
    shape_type: str                     # rect | accent_bar | divider | circle
    pos: tuple[float, float]           # (x, y)
    size: tuple[float, float]          # (w, h)
    fill_rgb: Optional[list[int]] = None
    line_rgb: Optional[list[int]] = None
    line_width: float = 1.0
    alpha: float = 1.0                  # 0-1 opacity for fill
    text: Optional[str] = None
    text_color_rgb: Optional[list[int]] = None
    text_font_size: float = 10.0


@dataclass
class ImageComponent:
    id: str
    image_path: str
    pos: tuple[float, float]
    size: tuple[float, float]
    fit: str = "contain"               # contain | cover | fill | stretch


@dataclass
class VizComponent:
    id: str
    viz_type: str                      # one of _RENDERERS keys
    viz_data: dict
    pos: tuple[float, float]
    size: tuple[float, float]
    caption: Optional[str] = None
    theme_name: str = "default"         # default | bosch
    badge: Optional[str] = None        # small tag shown top-right of viz area


from typing import Union
Component = Union[TextComponent, ShapeComponent, ImageComponent, VizComponent]


# ── Canvas ────────────────────────────────────────────────────────────────────

class Canvas:
    """
    Composable slide canvas. Units: inches, origin top-left.

    Coordinate system:
      x increases rightward
      y increases downward  (standard PPTX convention)

    Typical 16:9 slide: 13.333" wide × 7.5" tall

    Components are rendered in z-order (added order). Use z parameter
    to override the base z-order (default ordering: shape=0, text=1, image=2, viz=3).
    """

    # Default type ordering when z is not specified
    TYPE_ZORDER = {"shape": 0, "text": 1, "image": 2, "viz": 3}

    def __init__(
        self,
        width_in: float = 13.333,
        height_in: float = 7.5,
        theme_name: str = "default",
        accent_rgb: Optional[list[int]] = None,
        title_rgb: Optional[list[int]] = None,
        body_rgb: Optional[list[int]] = None,
        bg_rgb: Optional[list[int]] = None,
    ):
        self.width = width_in
        self.height = height_in
        self._components: list[tuple[int, Component]] = []   # (z, component)
        self._id_set: set[str] = set()
        self._viz_assets_dir: Optional[Path] = None
        self._theme_name = theme_name

        # Color tokens
        theme = BOSCH_THEME if theme_name == "bosch" else THEME
        self._accent = RGBColor(*(accent_rgb or [226, 0, 21] if theme_name == "bosch" else [37, 99, 235]))
        self._title  = RGBColor(*(title_rgb or [15, 23, 42]))
        self._body   = RGBColor(*(body_rgb  or [51, 65, 85]))
        self._bg     = RGBColor(*(bg_rgb or [255, 255, 255]))

    # ── Public API ────────────────────────────────────────────────────────────

    def add_text(
        self,
        id: str,
        content: str,
        pos: tuple[float, float],
        size: tuple[float, float],
        font_size: float = 12.0,
        bold: bool = False,
        color_rgb: Optional[list[int]] = None,
        alignment: str = "left",
        bullet: bool = False,
        level: int = 0,
        font_name: Optional[str] = None,
        wrap: bool = True,
        italic: bool = False,
        z: int = 1,
    ) -> None:
        """Add a text block."""
        self._register_id(id)
        c = TextComponent(
            id=id, content=content, pos=pos, size=size,
            font_size=font_size, bold=bold,
            color_rgb=color_rgb, alignment=alignment,
            bullet=bullet, level=level,
            font_name=font_name, wrap=wrap, italic=italic,
        )
        self._components.append((z + self.TYPE_ZORDER["text"], c))

    def add_viz(
        self,
        id: str,
        viz_type: str,
        viz_data: dict,
        pos: tuple[float, float],
        size: tuple[float, float],
        caption: Optional[str] = None,
        theme_name: Optional[str] = None,
        badge: Optional[str] = None,
        z: int = 3,
    ) -> None:
        """Add a visualization (renders to PNG via the appropriate renderer)."""
        self._register_id(id)
        c = VizComponent(
            id=id, viz_type=viz_type, viz_data=viz_data,
            pos=pos, size=size, caption=caption,
            theme_name=theme_name or self._theme_name, badge=badge,
        )
        self._components.append((z + self.TYPE_ZORDER["viz"], c))

    def add_image(
        self,
        id: str,
        image_path: str,
        pos: tuple[float, float],
        size: tuple[float, float],
        fit: str = "contain",
        z: int = 2,
    ) -> None:
        """Add an image block."""
        self._register_id(id)
        c = ImageComponent(id=id, image_path=image_path, pos=pos, size=size, fit=fit)
        self._components.append((z + self.TYPE_ZORDER["image"], c))

    def add_shape(
        self,
        id: str,
        shape_type: str,
        pos: tuple[float, float],
        size: tuple[float, float],
        fill_rgb: Optional[list[int]] = None,
        line_rgb: Optional[list[int]] = None,
        line_width: float = 1.0,
        alpha: float = 1.0,
        text: Optional[str] = None,
        text_color_rgb: Optional[list[int]] = None,
        text_font_size: float = 10.0,
        z: int = 0,
    ) -> None:
        """Add a shape. shape_type: rect | accent_bar | divider | circle | rounded_rect."""
        self._register_id(id)
        c = ShapeComponent(
            id=id, shape_type=shape_type, pos=pos, size=size,
            fill_rgb=fill_rgb, line_rgb=line_rgb,
            line_width=line_width, alpha=alpha,
            text=text, text_color_rgb=text_color_rgb, text_font_size=text_font_size,
        )
        self._components.append((z + self.TYPE_ZORDER["shape"], c))

    def add_badge(
        self,
        text: str,
        pos: tuple[float, float],
        color_rgb: Optional[list[int]] = None,
        text_color_rgb: Optional[list[int]] = None,
        font_size: float = 8.0,
        z: int = 5,
    ) -> None:
        """Convenience: add a pill badge (small rounded rect with centred text)."""
        w, h = max(len(text) * font_size * 0.018 + 0.2, 0.6), font_size * 0.06 + 0.15
        fill = color_rgb or [37, 99, 235]
        t_color = text_color_rgb or [255, 255, 255]
        self.add_shape(
            id=f"_badge_{text[:8]}",
            shape_type="rounded_rect",
            pos=pos, size=(w, h),
            fill_rgb=fill, line_rgb=None, line_width=0,
            text=text, text_color_rgb=t_color, text_font_size=font_size,
            z=z,
        )

    def query(self) -> dict[str, Any]:
        """
        Return the current canvas state — useful for the agent to decide
        what to place next.
        """
        used_area = []
        for z, c in sorted(self._components, key=lambda x: x[0]):
            if isinstance(c, TextComponent):
                used_area.append({
                    "id": c.id, "type": "text",
                    "pos": c.pos, "size": c.size,
                    "content_preview": c.content[:40],
                })
            elif isinstance(c, VizComponent):
                used_area.append({
                    "id": c.id, "type": "viz", "viz_type": c.viz_type,
                    "pos": c.pos, "size": c.size,
                })
            elif isinstance(c, ImageComponent):
                used_area.append({
                    "id": c.id, "type": "image",
                    "pos": c.pos, "size": c.size,
                })
            elif isinstance(c, ShapeComponent):
                used_area.append({
                    "id": c.id, "type": "shape", "shape_type": c.shape_type,
                    "pos": c.pos, "size": c.size,
                })

        # Compute free regions (very rough: just describe canvas bounds)
        return {
            "canvas_size": [self.width, self.height],
            "components": used_area,
            "component_count": len(self._components),
            "used_ids": sorted(self._id_set),
        }

    def to_json(self) -> dict:
        """Export canvas as serializable dict."""
        return {
            "width": self.width,
            "height": self.height,
            "theme": self._theme_name,
            "components": [
                self._component_to_dict(z, c)
                for z, c in sorted(self._components, key=lambda x: x[0])
            ],
        }

    def finish(self, output_path: str, assets_dir: Optional[str] = None) -> str:
        """
        Render all components to a PPTX file.

        Args:
            output_path:  .pptx output path
            assets_dir:   directory for rendered viz PNGs (default: <output>_assets/)

        Returns:
            str path to the output file
        """
        assets = Path(assets_dir) if assets_dir else Path(output_path).with_name(
            Path(output_path).stem + "_assets"
        )
        assets.mkdir(parents=True, exist_ok=True)
        self._viz_assets_dir = assets

        prs = Presentation()
        prs.slide_width  = Inches(self.width)
        prs.slide_height = Inches(self.height)
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)

        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.width), Inches(self.height),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self._bg
        bg_shape.line.fill.background()

        # Sort and render
        for z, c in sorted(self._components, key=lambda x: x[0]):
            self._render_component(slide, c)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    # ── Internal ──────────────────────────────────────────────────────────────

    def _register_id(self, id: str) -> None:
        if id in self._id_set:
            raise ValueError(f"Duplicate component id: '{id}'")
        self._id_set.add(id)

    def _component_to_dict(self, z: int, c: Component) -> dict[str, Any]:
        base = {"z": z, "id": c.id}
        if isinstance(c, TextComponent):
            return {**base, "type": "text",
                    "pos": c.pos, "size": c.size,
                    "content": c.content, "font_size": c.font_size,
                    "bold": c.bold, "color_rgb": c.color_rgb,
                    "alignment": c.alignment, "bullet": c.bullet, "level": c.level,
                    "italic": c.italic}
        elif isinstance(c, VizComponent):
            return {**base, "type": "viz", "viz_type": c.viz_type,
                    "viz_data": c.viz_data, "pos": c.pos, "size": c.size,
                    "caption": c.caption, "badge": c.badge}
        elif isinstance(c, ImageComponent):
            return {**base, "type": "image", "image_path": c.image_path,
                    "pos": c.pos, "size": c.size, "fit": c.fit}
        elif isinstance(c, ShapeComponent):
            return {**base, "type": "shape", "shape_type": c.shape_type,
                    "pos": c.pos, "size": c.size,
                    "fill_rgb": c.fill_rgb, "line_rgb": c.line_rgb,
                    "line_width": c.line_width, "alpha": c.alpha,
                    "text": c.text, "text_color_rgb": c.text_color_rgb,
                    "text_font_size": c.text_font_size}

    def _render_component(self, slide, c: Component) -> None:
        if isinstance(c, TextComponent):
            self._render_text(slide, c)
        elif isinstance(c, VizComponent):
            self._render_viz(slide, c)
        elif isinstance(c, ImageComponent):
            self._render_image(slide, c)
        elif isinstance(c, ShapeComponent):
            self._render_shape(slide, c)

    def _render_text(self, slide, c: TextComponent) -> None:
        from pptx.util import Pt as _Pt

        x, y = c.pos
        w, h = c.size
        left  = Inches(x)
        top   = Inches(y)
        width  = Inches(w)
        height = Inches(h)

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = c.wrap

        # Determine text color
        if c.color_rgb:
            rgb = RGBColor(c.color_rgb[0], c.color_rgb[1], c.color_rgb[2])
        elif c.bold:
            rgb = self._title
        else:
            rgb = self._body

        lines = c.content.split("\n")
        first = True
        for i, line in enumerate(lines):
            is_bullet = c.bullet or (c.content.startswith("•") or c.content.startswith("-"))
            line_text = line.strip()
            if not line_text:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(0)
                p.space_after  = Pt(0)
                continue

            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(0)
            p.space_after  = Pt(2)

            if c.bullet:
                prefix = "•  " if c.level == 0 else "◦  "
                # Build runs: prefix in accent, content in body
                run = p.add_run()
                run.text = prefix + line_text
            else:
                run = p.add_run()
                run.text = line_text

            run.font.size = Pt(c.font_size)
            run.font.bold = c.bold
            run.font.italic = c.italic
            run.font.color.rgb = rgb
            if c.font_name:
                try:
                    run.font.name = c.font_name
                except Exception:
                    pass

            # Alignment
            if c.alignment == "center":
                p.alignment = PP_ALIGN.CENTER
            elif c.alignment == "right":
                p.alignment = PP_ALIGN.RIGHT

    def _render_viz(self, slide, c: VizComponent) -> None:
        x, y = c.pos
        w, h = c.size
        left, top, width, height = Inches(x), Inches(y), Inches(w), Inches(h)

        renderer_key = c.viz_type.lower()
        if renderer_key not in _RENDERERS:
            raise ValueError(
                f"Unknown viz_type '{c.viz_type}'. "
                f"Available: {sorted(_RENDERERS.keys())}"
            )

        # Determine theme
        if c.theme_name == "bosch":
            from viz_theme import BOSCH_THEME as _T
        else:
            from viz_theme import THEME as _T

        # Build spec
        spec = copy.deepcopy(c.viz_data)
        if "title" not in spec and c.caption:
            spec["title"] = c.caption

        # Render PNG
        tmp_png = self._viz_assets_dir / f"canvas_{c.id}.png"
        tmp_pdf = self._viz_assets_dir / f"canvas_{c.id}.pdf"

        _RENDERERS[renderer_key](spec, str(tmp_png))

        # Fit image into box
        self._add_picture_fit(slide, str(tmp_png), x, y, w, h)

        # Caption below
        if c.caption:
            cap_tb = slide.shapes.add_textbox(
                Inches(x), top + height, width, Inches(0.35)
            )
            cap_tf = cap_tb.text_frame
            cap_p = cap_tf.paragraphs[0]
            cap_p.text = c.caption
            cap_p.font.size = Pt(9.0)
            cap_p.font.color.rgb = self._body
            cap_p.alignment = PP_ALIGN.LEFT

        # Badge top-right corner
        if c.badge:
            badge_w = len(c.badge) * 9 * 0.018 + 0.15
            badge_x = x + w - badge_w - 0.05
            badge_y = y + 0.05
            self._add_picture_fit(slide, str(tmp_png), badge_x, badge_y,
                                   badge_w, 0.28)  # placeholder; use shape instead
            # Use shape for badge
            bs = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(badge_x), Inches(badge_y), Inches(badge_w), Inches(0.28),
            )
            bs.fill.solid()
            bs.fill.fore_color.rgb = self._accent
            bs.line.fill.background()
            stb = bs.text_frame
            stp = stb.paragraphs[0]
            stp.text = c.badge
            stp.font.size = Pt(8)
            stp.font.color.rgb = RGBColor(255, 255, 255)
            stp.font.bold = True
            stp.alignment = PP_ALIGN.CENTER

    def _render_image(self, slide, c: ImageComponent) -> None:
        if not os.path.isfile(c.image_path):
            raise FileNotFoundError(f"Image not found: {c.image_path}")
        self._add_picture_fit(slide, c.image_path, c.pos[0], c.pos[1],
                               c.size[0], c.size[1])

    def _render_shape(self, slide, c: ShapeComponent) -> None:
        x, y = c.pos
        w, h = c.size
        left, top, width, height = Inches(x), Inches(y), Inches(w), Inches(h)

        # Select MSO_SHAPE
        shape_map = {
            "rect":         MSO_SHAPE.RECTANGLE,
            "accent_bar":   MSO_SHAPE.RECTANGLE,
            "divider":      MSO_SHAPE.RECTANGLE,
            "circle":       MSO_SHAPE.OVAL,
            "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        }
        mso = shape_map.get(c.shape_type, MSO_SHAPE.RECTANGLE)

        shp = slide.shapes.add_shape(mso, left, top, width, height)

        # Fill
        if c.fill_rgb is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(c.fill_rgb[0], c.fill_rgb[1], c.fill_rgb[2])
        else:
            shp.fill.background()

        # Line
        if c.line_rgb is not None:
            shp.line.color.rgb = RGBColor(c.line_rgb[0], c.line_rgb[1], c.line_rgb[2])
            shp.line.width = Pt(c.line_width)
        else:
            shp.line.fill.background()

        # Text
        if c.text:
            tf = shp.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = c.text
            p.font.size = Pt(c.text_font_size)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            if c.text_color_rgb:
                p.font.color.rgb = RGBColor(c.text_color_rgb[0], c.text_color_rgb[1], c.text_color_rgb[2])
            else:
                p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_picture_fit(
        self, slide, img_path: str,
        x: float, y: float, w: float, h: float,
    ) -> None:
        """Insert image letterboxed inside (x,y,w,h) box."""
        from PIL import Image as _PILImage
        try:
            with _PILImage.open(img_path) as im:
                img_w, img_h = im.size
        except Exception:
            return

        img_aspect = img_w / img_h
        box_aspect = w / h

        if img_aspect >= box_aspect:
            fit_w, fit_h = w, w / img_aspect
        else:
            fit_h, fit_w = h, h * img_aspect

        offset_left = (w - fit_w) / 2
        offset_top  = (h - fit_h) / 2

        slide.shapes.add_picture(
            img_path,
            Inches(x + offset_left),
            Inches(y + offset_top),
            width=Inches(fit_w),
            height=Inches(fit_h),
        )


# ── CLI ───────────────────────────────────────────────────────────────────────

def _canvas_from_json(data: dict) -> Canvas:
    """Reconstruct a Canvas from a serialised dict (to_json output)."""
    cv = Canvas(
        width_in=data.get("width", 13.333),
        height_in=data.get("height", 7.5),
        theme_name=data.get("theme", "default"),
    )
    zorder = {c["id"]: c.get("z", 0) for c in data.get("components", [])}
    type_order = {"shape": 0, "text": 1, "image": 2, "viz": 3}
    for comp in data.get("components", []):
        cid = comp["id"]
        ct  = comp["type"]
        base_z = zorder.get(cid, 0)
        ctype_z = type_order.get(ct, 1)
        z = comp.get("z", base_z + ctype_z)
        pos = tuple(comp["pos"])
        size = tuple(comp["size"])
        if ct == "text":
            cv.add_text(
                id=cid, content=comp.get("content", ""),
                pos=pos, size=size,
                font_size=comp.get("font_size", 12.0),
                bold=comp.get("bold", False),
                color_rgb=comp.get("color_rgb"),
                alignment=comp.get("alignment", "left"),
                bullet=comp.get("bullet", False),
                level=comp.get("level", 0),
                italic=comp.get("italic", False),
                z=z,
            )
        elif ct == "viz":
            cv.add_viz(
                id=cid, viz_type=comp.get("viz_type", ""),
                viz_data=comp.get("viz_data", {}),
                pos=pos, size=size,
                caption=comp.get("caption"),
                badge=comp.get("badge"),
                z=z,
            )
        elif ct == "image":
            cv.add_image(
                id=cid, image_path=comp.get("image_path", ""),
                pos=pos, size=size,
                fit=comp.get("fit", "contain"),
                z=z,
            )
        elif ct == "shape":
            cv.add_shape(
                id=cid, shape_type=comp.get("shape_type", "rect"),
                pos=pos, size=size,
                fill_rgb=comp.get("fill_rgb"),
                line_rgb=comp.get("line_rgb"),
                line_width=comp.get("line_width", 1.0),
                alpha=comp.get("alpha", 1.0),
                text=comp.get("text"),
                text_color_rgb=comp.get("text_color_rgb"),
                text_font_size=comp.get("text_font_size", 10.0),
                z=z,
            )
    return cv


def main() -> int:
    import argparse
    parser = argparse.ArgumentParser(description="Build a PPTX from a canvas JSON")
    parser.add_argument("--canvas", required=True, help="JSON file (canvas.to_json output)")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    parser.add_argument("--assets-dir", default=None, help="Viz assets directory")
    args = parser.parse_args()

    with open(args.canvas, encoding="utf-8") as f:
        data = json.load(f)

    cv = _canvas_from_json(data)
    out = cv.finish(args.output, assets_dir=args.assets_dir)
    print(f"Canvas saved: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
