#!/usr/bin/env python3
"""
Tesla FSD PPT Generator - Using Deployed Skills
直接使用已部署的 image-generation skill 生成 3 页 PPT
"""

import json
import sys
from pathlib import Path
from datetime import datetime


def main():
    print("=" * 80)
    print("Tesla FSD PPT Generation - Direct Execution")
    print("=" * 80)
    print()
    
    # 工作目录
    workspace = Path("/tmp/fsd_ppt_final/workspace")
    outputs = Path("/tmp/fsd_ppt_final/outputs")
    workspace.mkdir(parents=True, exist_ok=True)
    outputs.mkdir(parents=True, exist_ok=True)
    
    print(f"[1/4] Workspace: {workspace}")
    print(f"      Outputs:   {outputs}")
    print()
    
    # 准备 3 个幻灯片的提示词
    print("[2/4] Preparing slide prompts...")
    
    prompts = {
        "slide_1.json": {
            "prompt": "Tesla Full Self-Driving title slide. Bold red to blue gradient background (#E82127 to #0071E3). Center: large metallic silver 'FSD' text with glow. Neural network visualization with glowing nodes behind. Tesla Model 3 silhouette on right. Digital highway with autonomous path at bottom. Professional, futuristic, 16:9."
        },
        "slide_2.json": {
            "prompt": "Tech specifications slide for Tesla FSD. Center: Tesla Model 3 top-down view with sensor markings. Four floating glass cards around showing: 8 cameras + 12 sensors, real-time AI processing, 4+ billion fleet miles, safety redundancy. Premium tech aesthetic, gradient background."
        },
        "slide_3.json": {
            "prompt": "Future of mobility conclusion slide. Multiple Tesla vehicles on modern highway at golden sunset. Tech overlay elements, glowing effects. Bold 'The Future is Autonomous' text. Tesla logo with stats: 8M+ vehicles, 4B+ miles, 99.9% safety. Inspirational cinematography, premium quality."
        }
    }
    
    # 保存提示词文件
    for filename, content in prompts.items():
        prompt_file = workspace / filename
        with open(prompt_file, "w") as f:
            json.dump(content, f, indent=2)
    
    print(f"      ✓ {len(prompts)} prompts prepared")
    print()
    
    # 使用已部署的 image-generation 脚本
    print("[3/4] Generating slide images...")
    print("      Using: skills/public/image-generation/scripts/generate.py")
    print()
    
    # 调用已部署的脚本来生成图像
    import subprocess
    sys.path.insert(0, str(Path(__file__).parent / "backend" / "packages" / "harness"))
    
    skill_script = Path(__file__).parent / "skills/public/image-generation/scripts/generate.py"
    
    if not skill_script.exists():
        print(f"      ✗ Script not found: {skill_script}")
        return False
    
    # 为每个幻灯片生成图像
    for i, (filename, content) in enumerate(prompts.items(), 1):
        prompt_file = workspace / filename
        output_file = outputs / f"slide-{i:02d}.jpg"
        
        print(f"      Generating slide {i}/3...")
        
        try:
            result = subprocess.run(
                [
                    sys.executable,
                    str(skill_script),
                    "--prompt-file", str(prompt_file),
                    "--output-file", str(output_file),
                    "--aspect-ratio", "16:9"
                ],
                capture_output=True,
                text=True,
                timeout=600,
                cwd=str(skill_script.parent)
            )
            
            if result.returncode == 0:
                print(f"      ✓ slide-{i:02d}.jpg generated")
            else:
                print(f"      ✗ Error: {result.stderr}")
                return False
                
        except subprocess.TimeoutExpired:
            print(f"      ✗ Timeout generating slide {i}")
            return False
        except Exception as e:
            print(f"      ✗ Error: {e}")
            return False
    
    print()
    
    # 组合成 PPTX
    print("[4/4] Composing into PPTX...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        prs.title = "Tesla Full Self-Driving"
        
        # 添加幻灯片
        for i in range(1, 4):
            img_path = outputs / f"slide-{i:02d}.jpg"
            
            if not img_path.exists():
                print(f"      ✗ Image not found: {img_path}")
                return False
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            slide.shapes.add_picture(
                str(img_path),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            print(f"      ✓ Added slide {i}")
        
        # 保存
        output_pptx = outputs / "Tesla_FSD_3Slides.pptx"
        prs.save(str(output_pptx))
        
        print(f"      ✓ PPTX saved: {output_pptx}")
        
    except ImportError:
        print("      ✗ python-pptx not installed")
        print("      Run: pip install python-pptx")
        return False
    except Exception as e:
        print(f"      ✗ Error: {e}")
        return False
    
    print()
    print("=" * 80)
    print("✅ PPT Generation Complete!")
    print("=" * 80)
    print()
    print(f"Output: {output_pptx}")
    print(f"Size: {output_pptx.stat().st_size / (1024*1024):.2f} MB")
    print()
    print("Contents:")
    print("  • Slide 1: Title - Tesla FSD Brand")
    print("  • Slide 2: Technology - Core capabilities")
    print("  • Slide 3: Conclusion - Future vision")
    print()
    print(f"Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    return True


if __name__ == "__main__":
    try:
        success = main()
        sys.exit(0 if success else 1)
    except Exception as e:
        print(f"\n✗ Fatal error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
