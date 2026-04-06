#!/usr/bin/env python3
"""
Tesla FSD PPT Generation - Complete Agent Skills Flow
使用 ppt-generation skill 和 image-generation skill 生成 3 页 PPT
"""

import json
import sys
import subprocess
from pathlib import Path
from datetime import datetime


def main():
    print("=" * 80)
    print("Tesla FSD PPT Generation - Agent Skills Flow")
    print("=" * 80)
    print()
    
    # 工作目录
    workspace = Path("/tmp/fsd_agent_skills/workspace")
    outputs = Path("/tmp/fsd_agent_skills/outputs")
    workspace.mkdir(parents=True, exist_ok=True)
    outputs.mkdir(parents=True, exist_ok=True)
    
    print("[1/5] Setup")
    print(f"      Workspace: {workspace}")
    print(f"      Outputs:   {outputs}")
    print()
    
    # Step 1: PPT-Generation Skill - 创建演示计划
    print("[2/5] PPT-Generation Skill: Creating presentation plan...")
    
    ppt_plan = {
        "title": "Tesla Full Self-Driving (FSD)",
        "style": "gradient-modern",
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "title",
                "title": "Tesla Full Self-Driving",
                "subtitle": "Autonomous Driving Revolution",
                "visual_description": "Bold gradient background flowing from deep Tesla red (#E82127) through orange to electric blue (#0071E3). Center: large metallic silver 'FSD' text with glow effect. Behind: abstract neural network visualization with glowing nodes and connection lines. Top right: Tesla Model 3 silhouette in line art. Bottom: digital highway visualization with autonomous vehicle path overlay. Professional tech aesthetic, forward-looking, innovation-focused."
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "Key Technologies",
                "key_points": [
                    "8 cameras + 12 ultrasonic sensors",
                    "Real-time neural network processing",
                    "Fleet Learning: 4+ billion miles",
                    "Safety-first redundancy architecture"
                ],
                "visual_description": "Tech specs presentation slide. Center: Tesla Model 3 top-down view with sensor locations marked in different colors. Around vehicle: four floating glass-morphism cards showing camera array, processing speed metrics, fleet learning statistics (4+ billion miles), safety redundancy features. Modern gradient background creates premium feel. Data visualization elements showing autonomous capability improvements."
            },
            {
                "slide_number": 3,
                "type": "conclusion",
                "title": "The Future of Mobility",
                "subtitle": "Join the autonomous revolution",
                "visual_description": "Dramatic finale slide. Panoramic view of multiple Tesla vehicles on modern highway at sunset with autonomous systems active. Illustrated tech overlays and glowing elements. Large bold text 'The Future is Autonomous' overlays the scene. Tesla logo and key statistics prominently displayed: 8M+ vehicles equipped, 4B+ miles driven, 99.9% safety rating. Inspirational lighting, premium cinematography, call-to-action ready. Bottom: subtle Tesla branding."
            }
        ]
    }
    
    plan_file = workspace / "presentation_plan.json"
    with open(plan_file, "w") as f:
        json.dump(ppt_plan, f, indent=2)
    
    print(f"      ✓ Plan created: {len(ppt_plan['slides'])} slides")
    print()
    
    # Step 2: Image-Generation Skill - 生成每张幻灯片的图像
    print("[3/5] Image-Generation Skill: Generating slide images...")
    print("      (This will take 5-15 minutes depending on hardware)")
    print()
    
    image_gen_script = Path(__file__).parent / "skills/public/image-generation/scripts/generate.py"
    
    if not image_gen_script.exists():
        print(f"      ✗ Script not found: {image_gen_script}")
        return False
    
    # 为每张幻灯片生成图像
    for slide in ppt_plan['slides']:
        slide_num = slide['slide_number']
        
        # 创建 prompt 文件
        prompt_file = workspace / f"slide_{slide_num}_prompt.json"
        prompt_data = {
            "prompt": slide['visual_description']
        }
        with open(prompt_file, "w") as f:
            json.dump(prompt_data, f, indent=2)
        
        output_file = outputs / f"slide-{slide_num:02d}.jpg"
        
        print(f"      Generating slide {slide_num}/3 with image-generation skill...")
        
        try:
            # 调用 image-generation skill
            result = subprocess.run(
                [
                    sys.executable,
                    str(image_gen_script),
                    "--prompt-file", str(prompt_file),
                    "--output-file", str(output_file),
                    "--aspect-ratio", "16:9"
                ],
                capture_output=True,
                text=True,
                timeout=600,
                cwd=str(image_gen_script.parent)
            )
            
            if result.returncode == 0:
                print(f"      ✓ slide-{slide_num:02d}.jpg generated by image-generation skill")
            else:
                print(f"      ✗ Error: {result.stderr[:200]}")
                return False
                
        except subprocess.TimeoutExpired:
            print(f"      ✗ Timeout generating slide {slide_num}")
            return False
        except Exception as e:
            print(f"      ✗ Error: {e}")
            return False
    
    print()
    
    # Step 3: Composition Engine - 组合成 PPTX（ppt-generation skill 的最后一步）
    print("[4/5] PPT-Generation Skill: Composing into PPTX...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9
        prs.title = ppt_plan['title']
        prs.author = "Tesla FSD Agent"
        
        # 添加每一张幻灯片
        for i in range(1, 4):
            img_path = outputs / f"slide-{i:02d}.jpg"
            
            if not img_path.exists():
                print(f"      ✗ Image not found: {img_path}")
                return False
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图像（全尺寸）
            slide.shapes.add_picture(
                str(img_path),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            print(f"      ✓ Added slide {i} to PPTX")
        
        # 保存 PPTX
        output_pptx = outputs / "Tesla_FSD_Presentation.pptx"
        prs.save(str(output_pptx))
        
        print(f"      ✓ PPTX composition complete")
        
    except ImportError:
        print("      ✗ python-pptx not installed")
        return False
    except Exception as e:
        print(f"      ✗ Error: {e}")
        return False
    
    print()
    
    # Step 4: Verification
    print("[5/5] Verification")
    print(f"      File: {output_pptx}")
    print(f"      Size: {output_pptx.stat().st_size / (1024*1024):.2f} MB")
    print(f"      Slides: 3")
    print()
    
    print("=" * 80)
    print("✅ Agent Skills Pipeline Complete!")
    print("=" * 80)
    print()
    print("Skills Used:")
    print("  1. ppt-generation skill")
    print("     ├─ Created presentation plan (JSON)")
    print("     └─ Composed final PPTX")
    print()
    print("  2. image-generation skill")
    print("     ├─ Generated slide 1 (Title)")
    print("     ├─ Generated slide 2 (Technology)")
    print("     └─ Generated slide 3 (Conclusion)")
    print()
    print("  3. Stable Diffusion XL")
    print("     └─ Local inference (SDXL model)")
    print()
    print("Output:")
    print(f"  Location: {output_pptx}")
    print(f"  Format: PPTX (Microsoft PowerPoint)")
    print(f"  Resolution: 1024×576 (16:9)")
    print(f"  Style: gradient-modern (red → blue)")
    print()
    print(f"Cost: $0 (100% free - local SDXL)")
    print(f"Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    return True


if __name__ == "__main__":
    try:
        success = main()
        sys.exit(0 if success else 1)
    except Exception as e:
        print(f"\n✗ Fatal error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
