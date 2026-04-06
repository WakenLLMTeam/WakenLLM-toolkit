#!/usr/bin/env python3
"""
Tesla FSD PPT Generation - Complete Agent Skills Flow (Direct)
使用 ppt-generation skill 和 image-generation skill 生成 3 页 PPT
直接在主进程中调用，避免 subprocess 环境问题
"""

import json
import os
import sys
from pathlib import Path
from datetime import datetime


def _sdxl_local_files_only() -> bool:
    """Use only HF cache (no network). Set HF_HUB_OFFLINE=1 or SDXL_LOCAL_ONLY=1 when models are already downloaded."""
    return os.environ.get("HF_HUB_OFFLINE", "").strip().lower() in ("1", "true", "yes") or os.environ.get(
        "SDXL_LOCAL_ONLY", ""
    ).strip().lower() in ("1", "true", "yes")


def main():
    print("=" * 80)
    print("Tesla FSD PPT Generation - Agent Skills Flow (Direct)")
    print("=" * 80)
    print()
    
    # 工作目录
    workspace = Path("/tmp/fsd_agent_skills/workspace")
    outputs = Path("/tmp/fsd_agent_skills/outputs")
    workspace.mkdir(parents=True, exist_ok=True)
    outputs.mkdir(parents=True, exist_ok=True)
    
    print("[1/5] Setup")
    print(f"      Workspace: {workspace}")
    print(f"      Outputs:   {outputs}")
    print()
    
    # Step 1: PPT-Generation Skill - 创建演示计划
    print("[2/5] PPT-Generation Skill: Creating presentation plan...")
    
    ppt_plan = {
        "title": "Tesla Full Self-Driving (FSD)",
        "style": "gradient-modern",
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "title",
                "title": "Tesla Full Self-Driving",
                "subtitle": "Autonomous Driving Revolution",
                "visual_description": "Bold gradient background flowing from deep Tesla red (#E82127) through orange to electric blue (#0071E3). Center: large metallic silver 'FSD' text with glow effect. Behind: abstract neural network visualization with glowing nodes and connection lines. Top right: Tesla Model 3 silhouette in line art. Bottom: digital highway visualization with autonomous vehicle path overlay. Professional tech aesthetic, forward-looking, innovation-focused."
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "Key Technologies",
                "key_points": [
                    "8 cameras + 12 ultrasonic sensors",
                    "Real-time neural network processing",
                    "Fleet Learning: 4+ billion miles",
                    "Safety-first redundancy architecture"
                ],
                "visual_description": "Tech specs presentation slide. Center: Tesla Model 3 top-down view with sensor locations marked in different colors. Around vehicle: four floating glass-morphism cards showing camera array, processing speed metrics, fleet learning statistics (4+ billion miles), safety redundancy features. Modern gradient background creates premium feel. Data visualization elements showing autonomous capability improvements."
            },
            {
                "slide_number": 3,
                "type": "conclusion",
                "title": "The Future of Mobility",
                "subtitle": "Join the autonomous revolution",
                "visual_description": "Dramatic finale slide. Panoramic view of multiple Tesla vehicles on modern highway at sunset with autonomous systems active. Illustrated tech overlays and glowing elements. Large bold text 'The Future is Autonomous' overlays the scene. Tesla logo and key statistics prominently displayed: 8M+ vehicles equipped, 4B+ miles driven, 99.9% safety rating. Inspirational lighting, premium cinematography, call-to-action ready. Bottom: subtle Tesla branding."
            }
        ]
    }
    
    plan_file = workspace / "presentation_plan.json"
    with open(plan_file, "w") as f:
        json.dump(ppt_plan, f, indent=2)
    
    print(f"      ✓ Plan created: {len(ppt_plan['slides'])} slides")
    print()
    
    # Step 2: Image-Generation Skill - 生成每张幻灯片的图像
    print("[3/5] Image-Generation Skill: Generating slide images...")
    print("      (This will take 5-15 minutes depending on hardware)")
    print()
    
    # 直接导入并调用 generate.py 的函数
    try:
        print("      Loading SDXL model...")
        
        import torch
        from diffusers import StableDiffusionXLPipeline, DPMSolverMultistepScheduler
        from PIL import Image
        
        device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"      Using device: {device}")
        print()
        
        # 为每张幻灯片生成图像
        for slide in ppt_plan['slides']:
            slide_num = slide['slide_number']
            prompt = slide['visual_description']
            output_file = outputs / f"slide-{slide_num:02d}.jpg"
            
            print(f"      [{slide_num}/3] {slide['title'][:40]}...")
            
            # 加载模型（只在第一次加载）
            if slide_num == 1:
                _local = _sdxl_local_files_only()
                if _local:
                    print("            Loading SDXL from local Hugging Face cache only (no network)...")
                print(f"            Loading StableDiffusionXLPipeline (first load, ~5-10s)...")
                pipe = StableDiffusionXLPipeline.from_pretrained(
                    "stabilityai/stable-diffusion-xl-base-1.0",
                    torch_dtype=torch.float16 if device == "cuda" else torch.float32,
                    use_safetensors=True,
                    variant="fp16" if device == "cuda" else None,
                    local_files_only=_local,
                ).to(device)
                pipe.scheduler = DPMSolverMultistepScheduler.from_config(pipe.scheduler.config)
                if device == "cpu":
                    pipe.enable_sequential_cpu_offload()
                print(f"            ✓ Model loaded")
            
            # 生成图像
            print(f"            Generating image (inference: 20-60s)...")
            num_steps = 15 if device == "cpu" else 20
            
            with torch.inference_mode():
                image = pipe(
                    prompt=prompt,
                    height=576,
                    width=1024,
                    num_inference_steps=num_steps,
                    guidance_scale=7.5
                ).images[0]
            
            output_file.parent.mkdir(parents=True, exist_ok=True)
            image.save(str(output_file), quality=95)
            print(f"            ✓ Saved: {output_file.name} ({image.size[0]}×{image.size[1]})")
            print()
        
    except ImportError as e:
        print(f"      ✗ Import error: {e}")
        print(f"      Please ensure diffusers and torch are installed")
        return False
    except Exception as e:
        print(f"      ✗ Error during image generation: {e}")
        import traceback
        traceback.print_exc()
        return False
    
    print()
    
    # Step 3: Composition Engine - 组合成 PPTX（ppt-generation skill 的最后一步）
    print("[4/5] PPT-Generation Skill: Composing into PPTX...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9
        prs.title = ppt_plan['title']
        prs.author = "Tesla FSD Agent"
        
        # 添加每一张幻灯片
        for i in range(1, 4):
            img_path = outputs / f"slide-{i:02d}.jpg"
            
            if not img_path.exists():
                print(f"      ✗ Image not found: {img_path}")
                return False
            
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图像（全尺寸）
            slide.shapes.add_picture(
                str(img_path),
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            print(f"      ✓ Added slide {i} to PPTX")
        
        # 保存 PPTX
        output_pptx = outputs / "Tesla_FSD_Presentation.pptx"
        prs.save(str(output_pptx))
        
        print(f"      ✓ PPTX composition complete")
        
    except ImportError:
        print("      ✗ python-pptx not installed")
        return False
    except Exception as e:
        print(f"      ✗ Error: {e}")
        return False
    
    print()
    
    # Step 4: Verification
    print("[5/5] Verification")
    print(f"      File: {output_pptx}")
    print(f"      Size: {output_pptx.stat().st_size / (1024*1024):.2f} MB")
    print(f"      Slides: 3")
    print()
    
    print("=" * 80)
    print("✅ Agent Skills Pipeline Complete!")
    print("=" * 80)
    print()
    print("Skills Used:")
    print("  1. ppt-generation skill")
    print("     ├─ Created presentation plan (JSON)")
    print("     └─ Composed final PPTX")
    print()
    print("  2. image-generation skill")
    print("     ├─ Generated slide 1 (Title)")
    print("     ├─ Generated slide 2 (Technology)")
    print("     └─ Generated slide 3 (Conclusion)")
    print()
    print("  3. Stable Diffusion XL")
    print("     └─ Local inference (SDXL model)")
    print()
    print("Output:")
    print(f"  Location: {output_pptx}")
    print(f"  Format: PPTX (Microsoft PowerPoint)")
    print(f"  Resolution: 1024×576 (16:9)")
    print(f"  Style: gradient-modern (red → blue)")
    print()
    print(f"Cost: $0 (100% free - local SDXL)")
    print(f"Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print()
    print("=" * 80)
    
    return True


if __name__ == "__main__":
    try:
        success = main()
        sys.exit(0 if success else 1)
    except Exception as e:
        print(f"\n✗ Fatal error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
