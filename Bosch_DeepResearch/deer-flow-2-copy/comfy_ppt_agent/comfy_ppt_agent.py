#!/usr/bin/env python3
"""
ComfyUI 远程文生图 + 组装 PPTX（与本目录 ComfyUI 配套使用）。

前置:
  1) 在本目录执行 ./install_comfyui.sh 安装 ComfyUI 与 agent 依赖
  2) 将 Stable Diffusion 1.5 权重放到 ComfyUI/models/checkpoints/
     默认文件名: v1-5-pruned-emaonly.safetensors（可用环境变量 COMFY_CHECKPOINT 覆盖）
  3) 启动 ComfyUI: ./run_comfy_server.sh
  4) 生成: uv run python comfy_ppt_agent.py  或  python comfy_ppt_agent.py

环境变量:
  COMFY_SERVER   默认 127.0.0.1:8188
  COMFY_CHECKPOINT  checkpoints 下文件名，默认 v1-5-pruned-emaonly.safetensors
"""

from __future__ import annotations

import argparse
import io
import json
import os
import sys
import time
import uuid
from pathlib import Path
from typing import Any, Dict, List

import requests

# ---------------------------------------------------------------------------
# 路径（脚本与 ComfyUI 同位于 comfy_ppt_agent/）
# ---------------------------------------------------------------------------

ROOT = Path(__file__).resolve().parent
COMFY_ROOT = ROOT / "ComfyUI"
WORKFLOW_PATH = ROOT / "workflows" / "simple_t2i_api.json"
DEFAULT_OUT = Path.home() / ".cache" / "deerflow" / "comfy_ppt_agent_output"


def load_workflow_template() -> Dict[str, Any]:
    if not WORKFLOW_PATH.is_file():
        raise FileNotFoundError(f"缺少工作流文件: {WORKFLOW_PATH}")
    with open(WORKFLOW_PATH, encoding="utf-8") as f:
        return json.load(f)


def queue_prompt(server: str, workflow: Dict[str, Any], client_id: str) -> str:
    url = f"http://{server}/prompt"
    payload = {"prompt": workflow, "client_id": client_id}
    r = requests.post(url, json=payload, timeout=120)
    if r.status_code >= 400:
        raise RuntimeError(f"ComfyUI /prompt 失败: {r.status_code} {r.text}")
    data = r.json()
    return data["prompt_id"]


def wait_for_history(server: str, prompt_id: str, timeout_sec: float = 900.0) -> Dict[str, Any]:
    url = f"http://{server}/history/{prompt_id}"
    t0 = time.time()
    while time.time() - t0 < timeout_sec:
        r = requests.get(url, timeout=60)
        if r.status_code != 200:
            time.sleep(0.4)
            continue
        hist = r.json()
        if prompt_id in hist and hist[prompt_id].get("outputs"):
            return hist[prompt_id]
        time.sleep(0.4)
    raise TimeoutError(f"等待 ComfyUI 完成超时 ({timeout_sec}s): {prompt_id}")


def fetch_image(server: str, img: Dict[str, Any]) -> bytes:
    params = {
        "filename": img["filename"],
        "subfolder": img.get("subfolder", ""),
        "type": img.get("type", "output"),
    }
    r = requests.get(f"http://{server}/view", params=params, timeout=120)
    r.raise_for_status()
    return r.content


def build_workflow(
    template: Dict[str, Any],
    *,
    positive: str,
    negative: str,
    seed: int,
    width: int,
    height: int,
    ckpt_name: str,
    filename_prefix: str,
    steps: int = 20,
    cfg: float = 8.0,
) -> Dict[str, Any]:
    w = json.loads(json.dumps(template))
    w["3"]["inputs"]["seed"] = seed
    w["3"]["inputs"]["steps"] = int(steps)
    w["3"]["inputs"]["cfg"] = float(cfg)
    w["4"]["inputs"]["ckpt_name"] = ckpt_name
    w["5"]["inputs"]["width"] = int(width)
    w["5"]["inputs"]["height"] = int(height)
    w["6"]["inputs"]["text"] = positive
    w["7"]["inputs"]["text"] = negative
    w["9"]["inputs"]["filename_prefix"] = filename_prefix
    return w


def run_comfy_slide(
    server: str,
    template: Dict[str, Any],
    *,
    positive: str,
    negative: str,
    seed: int,
    ckpt: str,
    prefix: str,
    width: int,
    height: int,
) -> bytes:
    wf = build_workflow(
        template,
        positive=positive,
        negative=negative,
        seed=seed,
        width=width,
        height=height,
        ckpt_name=ckpt,
        filename_prefix=prefix,
    )
    cid = str(uuid.uuid4())
    pid = queue_prompt(server, wf, cid)
    hist = wait_for_history(server, pid)
    outs = hist.get("outputs") or {}
    for _node_id, node_out in outs.items():
        images = node_out.get("images") or []
        if images:
            return fetch_image(server, images[0])
    raise RuntimeError("ComfyUI 未返回图片输出，请检查 SaveImage 节点与日志")


def build_pptx(image_paths: List[Path], title: str, out_pptx: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    try:
        prs.core_properties.title = title
        prs.core_properties.author = "comfy_ppt_agent"
    except Exception:
        pass
    blank = prs.slide_layouts[6]
    for p in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    prs.save(str(out_pptx))


def default_demo_slides() -> List[Dict[str, str]]:
    return [
        {
            "title": "Slide 1",
            "prompt": "masterpiece, abstract logic diagram, Venn circles, soft blue gradient, clean infographic, no text",
        },
        {
            "title": "Slide 2",
            "prompt": "masterpiece, flowchart arrows nodes, navy and gold, professional presentation, no readable text",
        },
        {
            "title": "Slide 3",
            "prompt": "masterpiece, futuristic data visualization, purple cyan glow, dark background, no text",
        },
    ]


def main() -> int:
    ap = argparse.ArgumentParser(description="ComfyUI 文生图 + PPTX 组装（需先启动 ComfyUI 服务）")
    ap.add_argument("--server", default=os.environ.get("COMFY_SERVER", "127.0.0.1:8188"), help="ComfyUI 地址")
    ap.add_argument("--checkpoint", default=os.environ.get("COMFY_CHECKPOINT", "v1-5-pruned-emaonly.safetensors"))
    ap.add_argument("--out-dir", type=Path, default=DEFAULT_OUT)
    ap.add_argument("--pptx-name", default="ComfyUI_Generated.pptx")
    ap.add_argument("--width", type=int, default=1024)
    ap.add_argument("--height", type=int, default=576)
    ap.add_argument("--steps", type=int, default=20)
    ap.add_argument("--negative", default="worst quality, blurry, watermark, ugly, deformed, readable text, letters")
    ap.add_argument("--dry-run", action="store_true", help="只打印计划，不调用 ComfyUI")
    args = ap.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    template = load_workflow_template()

    slides = default_demo_slides()
    print(f"输出目录: {args.out_dir}")
    print(f"ComfyUI: http://{args.server}  checkpoint={args.checkpoint}")

    saved: List[Path] = []
    for i, s in enumerate(slides, start=1):
        prefix = f"ppt_{i:02d}"
        out_jpg = args.out_dir / f"slide-{i:02d}.jpg"
        print(f"[{i}/{len(slides)}] {s['title']}: {s['prompt'][:80]}...")
        if args.dry_run:
            continue
        seed = 100000 + i * 1337
        png_bytes = run_comfy_slide(
            args.server,
            template,
            positive=s["prompt"],
            negative=args.negative,
            seed=seed,
            ckpt=args.checkpoint,
            prefix=prefix,
            width=args.width,
            height=args.height,
        )
        from PIL import Image

        Image.open(io.BytesIO(png_bytes)).convert("RGB").save(str(out_jpg), quality=95)
        saved.append(out_jpg)
        print(f"  -> {out_jpg}")

    if args.dry_run:
        print("dry-run: 未生成文件")
        return 0

    if not saved:
        return 1

    pptx_path = args.out_dir / args.pptx_name
    build_pptx(saved, "ComfyUI PPT Demo", pptx_path)
    plan = {"slides": slides, "server": args.server, "checkpoint": args.checkpoint}
    with open(args.out_dir / "plan.json", "w", encoding="utf-8") as f:
        json.dump(plan, f, indent=2, ensure_ascii=False)
    print(f"PPTX: {pptx_path}")
    return 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except KeyboardInterrupt:
        print("中断", file=sys.stderr)
        sys.exit(130)
