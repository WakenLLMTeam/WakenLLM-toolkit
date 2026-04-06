#!/usr/bin/env python3
"""
PPT Generation Runner - 改进的脚本执行器
- 处理容器路径 vs 本地路径映射
- 统一环境设置
- 提供详细日志和进度跟踪
- 支持 SDXL 本地推理（diffusers）或 ComfyUI HTTP API（推荐：与仓库内 comfy_ppt_agent/ComfyUI 配合）

运行:
  uv run python ../scripts/ppt_generation_runner.py                      # 默认: 4 套博世主题 + 归档 12 组 prompt
  uv run python ../scripts/ppt_generation_runner.py --mode tesla       # 仅原版 Tesla FSD 单套
  uv run python ../scripts/ppt_generation_runner.py --mode illustrated # 1 套 3 页：左图右文（中文）
  uv run python ../scripts/ppt_generation_runner.py --mode logic       # 1 套 3 页：逻辑推理可视化（全图）

ComfyUI 模式（需先启动 ComfyUI 服务，并将 SDXL base 权重放入 ComfyUI/models/checkpoints/）:
  ./scripts/run_comfyui_server.sh    # 或: cd comfy_ppt_agent/ComfyUI && python main.py --listen 127.0.0.1 --port 8188
  uv run python ../scripts/ppt_generation_runner.py --backend comfy --mode illustrated
  环境变量: DEERFLOW_COMFY_HOST, DEERFLOW_COMFY_PORT, DEERFLOW_COMFY_CHECKPOINT（默认 sd_xl_base_1.0.safetensors）
"""

import argparse
import io
import json
import logging
import os
import random
import re
import sys
import time
import urllib.error
import urllib.parse
import urllib.request
import uuid
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, Any, List, Union

from tqdm.auto import tqdm


# ============================================================================
# 日志配置
# ============================================================================

def setup_logging():
    """配置详细的日志系统"""
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s [%(levelname)s] %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    return logging.getLogger(__name__)


logger = setup_logging()


# ============================================================================
# 多份演示素材（车辆 / 自动驾驶 / 博世风格）；每份 3 页；提示词保持简短以利 CLIP
# 视觉为 AI 生成示意，非官方物料。博世常见品牌色：红 #E30613、白、深灰。
# ============================================================================

def _slug(s: str) -> str:
    s = re.sub(r"[^\w\s-]", "", s, flags=re.ASCII)
    s = re.sub(r"[-\s]+", "_", s.strip()).lower()
    return s[:48] or "deck"


BATCH_DECKS: List[Dict[str, Any]] = [
    {
        "id": "bosch_mobility_adas",
        "pptx_name": "Bosch_Mobility_ADAS.pptx",
        "plan": {
            "title": "Bosch Mobility & ADAS",
            "style": "corporate-red-gray",
            "aspect_ratio": "16:9",
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "title": "Bosch Mobility Solutions",
                    "subtitle": "Safer roads, smarter vehicles",
                    "visual_description": "Title slide, Bosch red and white corporate style, subtle gray grid. Silhouette of modern car, abstract ADAS sensor arcs, clean minimalist keynote, professional automotive supplier aesthetic.",
                },
                {
                    "slide_number": 2,
                    "type": "content",
                    "title": "ADAS Building Blocks",
                    "key_points": ["Sensors", "Compute", "Actuation", "Software"],
                    "visual_description": "Infographic layout, vehicle top view, camera radar lidar icons around car, central ECU block, blue and red accents on light gray, technical blueprint style, no readable text.",
                },
                {
                    "slide_number": 3,
                    "type": "conclusion",
                    "title": "Driving Assistance Today",
                    "subtitle": "From assistance to automation",
                    "visual_description": "Sunset highway, multiple cars with soft glow trails, futuristic but realistic, warm light, sense of motion, wide cinematic 16:9, premium mobility closing slide.",
                },
            ],
        },
    },
    {
        "id": "autonomous_driving_stack",
        "pptx_name": "Autonomous_Driving_Stack.pptx",
        "plan": {
            "title": "Autonomous Driving: Perception to Action",
            "style": "tech-dark",
            "aspect_ratio": "16:9",
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "title": "Autonomous Driving Stack",
                    "subtitle": "Perception, planning, control",
                    "visual_description": "Dark tech background, neon cyan lines, self-driving car wireframe center, lidar point cloud hint, sci-fi dashboard mood, high contrast, no logos.",
                },
                {
                    "slide_number": 2,
                    "type": "content",
                    "title": "Sense, Decide, Act",
                    "key_points": ["Perception", "Localization", "Planning"],
                    "visual_description": "Three horizontal panels, icons for eyes brain steering, flow arrows left to right, dark blue UI, glass cards, abstract vehicle in middle, enterprise presentation style.",
                },
                {
                    "slide_number": 3,
                    "type": "conclusion",
                    "title": "Safety First",
                    "subtitle": "Validation and redundancy",
                    "visual_description": "Shield and checkmark motif, layered safety rings around car diagram, calm green and blue tones, trustworthy engineering finale slide.",
                },
            ],
        },
    },
    {
        "id": "connected_vehicle_v2x",
        "pptx_name": "Connected_Vehicle_V2X.pptx",
        "plan": {
            "title": "Connected Vehicle & V2X",
            "style": "network-blue",
            "aspect_ratio": "16:9",
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "title": "Connected Mobility",
                    "subtitle": "Vehicle to everything",
                    "visual_description": "Smart city aerial view soft focus, cars linked by glowing network lines, IoT nodes, blue and white palette, modern telecom style illustration.",
                },
                {
                    "slide_number": 2,
                    "type": "content",
                    "title": "V2X Ecosystem",
                    "key_points": ["V2V", "V2I", "Cloud"],
                    "visual_description": "Diagram style, car talking to traffic light and cloud, dotted wireless links, isometric simple shapes, pastel blues, clean infographic without text.",
                },
                {
                    "slide_number": 3,
                    "type": "conclusion",
                    "title": "Future of Traffic Flow",
                    "subtitle": "Cooperative systems",
                    "visual_description": "Smooth multi-lane traffic from above, cooperative merge visualization, golden hour light, optimistic sustainable mobility theme.",
                },
            ],
        },
    },
    {
        "id": "bosch_sensors_automotive",
        "pptx_name": "Bosch_Automotive_Sensors.pptx",
        "plan": {
            "title": "Automotive Sensing & Radar",
            "style": "engineering",
            "aspect_ratio": "16:9",
            "slides": [
                {
                    "slide_number": 1,
                    "type": "title",
                    "title": "Precision Sensing",
                    "subtitle": "Radar, camera, ultrasonics",
                    "visual_description": "Engineering title slide, Bosch-like red stripe accent, exploded sensor kit around vehicle bumper, technical cross-section mood, white gray background.",
                },
                {
                    "slide_number": 2,
                    "type": "content",
                    "title": "Sensor Fusion",
                    "key_points": ["Radar", "Vision", "USS"],
                    "visual_description": "Sensor fusion concept, overlapping transparent cones from car front, heatmap style, data fusion block merging streams, cool technical illustration.",
                },
                {
                    "slide_number": 3,
                    "type": "conclusion",
                    "title": "Reliability in Motion",
                    "subtitle": "Automotive grade quality",
                    "visual_description": "Vehicle on test track dusk, calibration targets blurred background, professional OEM supplier closing visual, subtle motion blur.",
                },
            ],
        },
    },
]

# 图文并茂：SDXL 仅生成左侧配图（prompt 要求无文字）；右侧中文由 PPT 叠加，保证可读。
ILLUSTRATED_AD_DECK: Dict[str, Any] = {
    "id": "ad_illustrated_zh",
    "pptx_name": "AD_Illustrated_Text_ZH.pptx",
    "plan": {
        "title": "自动驾驶：图文并茂导读",
        "style": "split-image-text",
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "content",
                "title": "Overview",
                "visual_description": "Cinematic highway at blue hour, single modern sedan, subtle sensor glow rings, clean automotive keynote, wide shot, no text no letters.",
                "overlay_title": "自动驾驶：从辅助到协同",
                "overlay_body": "• 多传感器融合感知周围环境\n• 人机共驾与责任边界\n• 面向量产的功能安全与验证",
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "Perception stack",
                "visual_description": "Dark blue background, stylized lidar point cloud around vehicle silhouette, soft neural glow, engineering illustration, no text no logos.",
                "overlay_title": "感知与决策链",
                "overlay_body": "• 视觉 / 毫米波雷达 / 激光协同\n• 定位、预测与路径规划\n• 控制闭环与冗余架构",
            },
            {
                "slide_number": 3,
                "type": "content",
                "title": "V2X future",
                "visual_description": "Bird eye smart intersection, connected vehicles light trails, roadside unit hints, green blue palette, illustration, no readable text.",
                "overlay_title": "车路协同与未来出行",
                "overlay_body": "• V2X 与交通流协同优化\n• 数据链路与功能安全\n• 可持续的智慧交通愿景",
            },
        ],
    },
}

# 逻辑推理可视化：文氏图、三段论链、真值表/门电路意象；prompt 避免可读文字（防 AI 乱码）
LOGIC_REASONING_DECK: Dict[str, Any] = {
    "id": "logic_reasoning_viz",
    "pptx_name": "Logic_Reasoning_Visualization.pptx",
    "plan": {
        "title": "Logical Reasoning — Visual Maps",
        "style": "logic-edu",
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "content",
                "title": "Sets and Venn",
                "visual_description": "Three overlapping Venn circles, soft blue teal coral, white background, abstract diagram, no text.",
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "Syllogism flow",
                "visual_description": "Vertical flowchart three round nodes linked by arrows, navy and gold, clean lines, no text.",
            },
            {
                "slide_number": 3,
                "type": "content",
                "title": "Truth and gates",
                "visual_description": "Neon grid like truth table, simple AND OR NOT icons, purple cyan on dark slate, no text.",
            },
        ],
    },
}

TESLA_LEGACY_PLAN: Dict[str, Any] = {
    "title": "Tesla Full Self-Driving (FSD)",
    "style": "gradient-modern",
    "aspect_ratio": "16:9",
    "slides": [
        {
            "slide_number": 1,
            "type": "title",
            "title": "Tesla Full Self-Driving",
            "subtitle": "Autonomous Driving Revolution",
            "visual_description": "Bold gradient background flowing from deep Tesla red (#E82127) through orange to electric blue (#0071E3). Center: large metallic silver 'FSD' text with glow effect. Behind: abstract neural network visualization with glowing nodes and connection lines. Top right: Tesla Model 3 silhouette in line art. Bottom: digital highway visualization with autonomous vehicle path overlay. Professional tech aesthetic, forward-looking, innovation-focused.",
        },
        {
            "slide_number": 2,
            "type": "content",
            "title": "Key Technologies",
            "key_points": [
                "8 cameras + 12 ultrasonic sensors",
                "Real-time neural network processing",
                "Fleet Learning: 4+ billion miles",
                "Safety-first redundancy architecture",
            ],
            "visual_description": "Tech specs presentation slide. Center: Tesla Model 3 top-down view with sensor locations marked in different colors. Around vehicle: four floating glass-morphism cards showing camera array, processing speed metrics, fleet learning statistics (4+ billion miles), safety redundancy features. Modern gradient background creates premium feel. Data visualization elements showing autonomous capability improvements.",
        },
        {
            "slide_number": 3,
            "type": "conclusion",
            "title": "The Future of Mobility",
            "subtitle": "Join the autonomous revolution",
            "visual_description": "Dramatic finale slide. Panoramic view of multiple Tesla vehicles on modern highway at sunset with autonomous systems active. Illustrated tech overlays and glowing elements. Large bold text 'The Future is Autonomous' overlays the scene. Tesla logo and key statistics prominently displayed: 8M+ vehicles equipped, 4B+ miles driven, 99.9% safety rating. Inspirational lighting, premium cinematography, call-to-action ready. Bottom: subtle Tesla branding.",
        },
    ],
}


# ============================================================================
# 路径映射
# ============================================================================

class PathMapper:
    """处理容器路径和本地路径的映射"""
    
    def __init__(self, project_root: Optional[Path] = None):
        if project_root is None:
            # 自动检测项目根路径
            self.project_root = Path(__file__).parent.parent
        else:
            self.project_root = project_root
        
        self.skills_dir = self.project_root / "skills"
        self.cache_dir = Path.home() / ".cache" / "deerflow"
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"项目根路径: {self.project_root}")
        logger.info(f"Skills 目录: {self.skills_dir}")
        logger.info(f"缓存目录: {self.cache_dir}")
    
    def resolve_path(self, container_path: str) -> Path:
        """
        将容器路径或本地路径转换为绝对路径
        
        Examples:
            /mnt/skills/public/image-generation/SKILL.md
            → {project_root}/skills/public/image-generation/SKILL.md
            
            /mnt/user-data/workspace/prompt.json
            → {cache_dir}/workspace/prompt.json
        """
        p = Path(container_path)
        
        # 容器 skills 路径
        if str(p).startswith("/mnt/skills/"):
            rel_path = p.relative_to("/mnt/skills/")
            resolved = self.skills_dir / rel_path
            logger.debug(f"映射 {container_path} → {resolved}")
            return resolved
        
        # 容器 user-data 路径
        if str(p).startswith("/mnt/user-data/"):
            rel_path = p.relative_to("/mnt/user-data/")
            resolved = self.cache_dir / rel_path
            resolved.parent.mkdir(parents=True, exist_ok=True)
            logger.debug(f"映射 {container_path} → {resolved}")
            return resolved
        
        # 本地路径直接使用
        if p.is_absolute():
            logger.debug(f"使用本地路径: {p}")
            return p
        
        # 相对路径
        resolved = self.project_root / p
        logger.debug(f"映射 {container_path} → {resolved}")
        return resolved


# ============================================================================
# 模型管理
# ============================================================================

class ModelManager:
    """SDXL 模型管理"""
    
    def __init__(self):
        self.model_id = "stabilityai/stable-diffusion-xl-base-1.0"
        self.cache_dir = Path.home() / ".cache" / "huggingface" / "hub"
        self.pipe = None
        self.device = None
    
    def check_model_available(self) -> bool:
        """检查模型是否已下载"""
        if not self.cache_dir.exists():
            logger.warning(f"HF 缓存目录不存在: {self.cache_dir}")
            return False
        
        # 检查模型文件夹
        model_folders = list(self.cache_dir.glob("models--*"))
        if model_folders:
            total_size = sum(f.stat().st_size for f in model_folders[0].rglob("*") if f.is_file())
            size_gb = total_size / 1e9
            logger.info(f"✓ 模型已缓存 ({size_gb:.1f} GB)")
            return True
        
        return False
    
    def load_model(self):
        """加载 SDXL 模型"""
        import torch
        from diffusers import StableDiffusionXLPipeline, DPMSolverMultistepScheduler
        
        logger.info("=" * 70)
        logger.info("加载 SDXL 模型...")
        logger.info("=" * 70)
        
        # 检测设备
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"使用设备: {self.device}")
        
        # 检查模型缓存
        self.check_model_available()
        
        # 加载模型
        start_time = time.time()
        
        logger.info(f"正在从 Hugging Face 加载 {self.model_id}...")
        logger.info("(第一次使用会下载模型，请耐心等待...)")
        
        try:
            # HF 权重在 ~/.cache/huggingface/hub/models--*；勿用 hub 的父目录，否则会落到错误的 models-- 路径
            # local_files_only：避免每次向 huggingface.co 校验元数据导致网络超时（模型已离线可用时）
            self.pipe = StableDiffusionXLPipeline.from_pretrained(
                self.model_id,
                torch_dtype=torch.float16 if self.device == "cuda" else torch.float32,
                use_safetensors=True,
                variant="fp16" if self.device == "cuda" else None,
                cache_dir=str(self.cache_dir),
                local_files_only=True,
            ).to(self.device)
            
            # 配置调度器
            self.pipe.scheduler = DPMSolverMultistepScheduler.from_config(
                self.pipe.scheduler.config
            )
            
            # CPU 优化
            if self.device == "cpu":
                logger.info("启用 CPU 内存优化...")
                self.pipe.enable_sequential_cpu_offload()

            # 使用脚本内 tqdm 显示去噪进度，避免与 logging 双进度条混淆
            self.pipe.set_progress_bar_config(disable=True)
            
            elapsed = time.time() - start_time
            logger.info(f"✓ 模型加载完成 ({elapsed:.1f}s)")
            logger.info("=" * 70)
            
            return self.pipe
            
        except Exception as e:
            logger.error(f"模型加载失败: {e}")
            raise

    def _clip_safe_prompt(self, prompt: str) -> str:
        """SDXL 第一路 CLIP 最多 77 token；过长会触发 transformers 告警并丢弃尾部。"""
        if self.pipe is None:
            return prompt
        tok = self.pipe.tokenizer
        max_len = min(77, int(getattr(tok, "model_max_length", 77) or 77))
        enc = tok(prompt, truncation=True, max_length=max_len, return_tensors="pt")
        out = tok.batch_decode(enc.input_ids, skip_special_tokens=True)[0].strip()
        if out != prompt.strip():
            logger.info("提示词已按 CLIP %s token 上限截断（避免尾部被静默丢弃）", max_len)
        return out
    
    def generate_image(
        self,
        prompt: str,
        width: int = 1024,
        height: int = 576,
        *,
        slide_title: Optional[str] = None,
    ) -> Any:
        """生成图像；去噪过程用 tqdm 显示每步进度。"""
        import torch

        if self.pipe is None:
            raise RuntimeError("模型未加载")

        prompt = self._clip_safe_prompt(prompt)

        num_steps = 15 if self.device == "cpu" else 20
        title_short = (slide_title or "slide")[:40]
        logger.info(f"生成图像 [{title_short}] 步数={num_steps} 设备={self.device}")
        logger.debug(f"prompt: {prompt[:120]}...")

        pbar = tqdm(
            total=num_steps,
            desc=f"SDXL 去噪 · {title_short}",
            unit="step",
            leave=False,
            dynamic_ncols=True,
            mininterval=0.2,
            position=1,
        )

        def _on_step_end(_pipe, _step_index, _timestep, callback_kwargs):
            pbar.update(1)
            return callback_kwargs

        start_time = time.time()
        try:
            with torch.inference_mode():
                image = self.pipe(
                    prompt=prompt,
                    height=height,
                    width=width,
                    num_inference_steps=num_steps,
                    guidance_scale=7.5,
                    callback_on_step_end=_on_step_end,
                ).images[0]
        except Exception as e:
            logger.error(f"生成失败: {e}")
            raise
        finally:
            pbar.close()

        elapsed = time.time() - start_time
        logger.info(f"✓ 生成完成 ({elapsed:.1f}s)")

        return image


# ============================================================================
# PPT 生成主函数
# ============================================================================


def run_one_deck(
    ppt_plan: Dict[str, Any],
    deck_dir: Path,
    pptx_filename: str,
    model_manager: ModelManager,
    *,
    deck_label: str = "",
) -> bool:
    """
    在已加载的 model_manager 上生成一套 3 页幻灯片并写出 pptx。
    deck_dir 内写入 slide-01.jpg … 与 presentation_plan.json。
    """
    deck_dir.mkdir(parents=True, exist_ok=True)
    plan_file = deck_dir / "presentation_plan.json"
    with open(plan_file, "w", encoding="utf-8") as f:
        json.dump(ppt_plan, f, indent=2, ensure_ascii=False)

    prefix = f"[{deck_label}] " if deck_label else ""
    logger.info("%s✓ 演示计划: %s | %s", prefix, ppt_plan["title"], plan_file)

    slides = ppt_plan["slides"]
    n_slides = len(slides)
    logger.info("=" * 70)
    logger.info("%s[3/4] 生成幻灯片图像 (%s 张) → %s", prefix, n_slides, deck_dir.name)
    logger.info("=" * 70)

    generated_images: List[str] = []

    try:
        slide_pbar = tqdm(
            slides,
            desc=f"{deck_label or 'deck'} · 幻灯片",
            unit="张",
            dynamic_ncols=True,
            leave=True,
            position=0,
        )
        for slide in slide_pbar:
            slide_num = slide["slide_number"]
            slide_pbar.set_postfix_str(slide["title"][:28] + ("…" if len(slide["title"]) > 28 else ""))
            prompt = slide["visual_description"]
            output_file = deck_dir / f"slide-{slide_num:02d}.jpg"
            image = model_manager.generate_image(prompt, slide_title=slide["title"])
            image.save(str(output_file), quality=95)
            logger.info("%s✓ 已保存: %s", prefix, output_file.name)
            generated_images.append(str(output_file))
        logger.info("")
    except Exception as e:
        logger.error("%s图像生成失败: %s", prefix, e)
        import traceback
        traceback.print_exc()
        return False

    logger.info("=" * 70)
    logger.info("%s[4/4] 组合 PPTX", prefix)
    logger.info("=" * 70)

    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        prs.title = ppt_plan["title"]
        prs.author = "DeerFlow PPT Generator"

        for _i, img_path in enumerate(
            tqdm(generated_images, desc=f"{deck_label or 'deck'} · PPTX", unit="页", dynamic_ncols=True),
            start=1,
        ):
            blank_slide_layout = prs.slide_layouts[6]
            sld = prs.slides.add_slide(blank_slide_layout)
            sld.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        output_pptx = deck_dir / pptx_filename
        prs.save(str(output_pptx))
        logger.info("%s✓ PPTX: %s (%s 页, %.1f MB)", prefix, output_pptx.name, len(prs.slides), output_pptx.stat().st_size / (1024 * 1024))
        logger.info("")
    except Exception as e:
        logger.error("%sPPTX 组合失败: %s", prefix, e)
        import traceback
        traceback.print_exc()
        return False

    return True


def export_bosch_batch_12_prompts(workspace: Path, script_dir: Path) -> Path:
    """将 4 套×3 页共 12 组 visual_description 写入 JSON（与代码中 BATCH_DECKS 一致）。"""
    payload: Dict[str, Any] = {
        "exported_at": datetime.now().isoformat(timespec="seconds"),
        "note": "博世批量四套演示，每页一条 visual_description（SDXL 英文提示词）",
        "deck_count": len(BATCH_DECKS),
        "slides_per_deck": 3,
        "decks": [],
    }
    for d in BATCH_DECKS:
        slides_out = []
        for s in d["plan"]["slides"]:
            slides_out.append(
                {
                    "slide_number": s["slide_number"],
                    "slide_title": s["title"],
                    "visual_description": s["visual_description"],
                }
            )
        payload["decks"].append(
            {
                "deck_id": d["id"],
                "pptx_name": d["pptx_name"],
                "presentation_title": d["plan"]["title"],
                "slides": slides_out,
            }
        )
    out_ws = workspace / "bosch_batch_12_prompts.json"
    out_script = script_dir / "bosch_batch_12_prompts.json"
    for out in (out_ws, out_script):
        out.parent.mkdir(parents=True, exist_ok=True)
        with open(out, "w", encoding="utf-8") as f:
            json.dump(payload, f, indent=2, ensure_ascii=False)
    logger.info("已归档 12 组 prompt → %s 与 %s", out_ws, out_script)
    return out_script


def run_one_deck_illustrated(
    deck: Dict[str, Any],
    deck_dir: Path,
    model_manager: ModelManager,
    *,
    deck_label: str = "",
) -> bool:
    """左图右文：生成 3 张配图 + 带中文矢量标题/正文的 pptx。"""
    ppt_plan = deck["plan"]
    pptx_filename = deck["pptx_name"]
    deck_dir.mkdir(parents=True, exist_ok=True)
    plan_file = deck_dir / "presentation_plan.json"
    with open(plan_file, "w", encoding="utf-8") as f:
        json.dump(ppt_plan, f, indent=2, ensure_ascii=False)

    prefix = f"[{deck_label}] " if deck_label else ""
    logger.info("%s图文并茂套: %s | %s", prefix, ppt_plan["title"], plan_file)

    slides = ppt_plan["slides"]
    generated_images: List[str] = []

    try:
        slide_pbar = tqdm(slides, desc=f"{deck_label or 'ill'} · 幻灯片", unit="张", dynamic_ncols=True, leave=True, position=0)
        for slide in slide_pbar:
            sn = slide["slide_number"]
            slide_pbar.set_postfix_str((slide.get("overlay_title") or slide["title"])[:20])
            out_jpg = deck_dir / f"slide-{sn:02d}.jpg"
            image = model_manager.generate_image(slide["visual_description"], slide_title=slide["title"])
            image.save(str(out_jpg), quality=95)
            logger.info("%s✓ 已保存: %s", prefix, out_jpg.name)
            generated_images.append(str(out_jpg))
        logger.info("")
    except Exception as e:
        logger.error("%s图像生成失败: %s", prefix, e)
        import traceback
        traceback.print_exc()
        return False

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        try:
            prs.core_properties.title = ppt_plan["title"]
            prs.core_properties.author = "DeerFlow PPT Generator"
        except Exception:
            pass

        img_w = Inches(5.85)
        full_h = Inches(5.625)
        panel_left = Inches(5.85)
        panel_w = Inches(4.15)

        for idx, img_path in enumerate(generated_images):
            slide_meta = slides[idx]
            blank = prs.slide_layouts[6]
            sld = prs.slides.add_slide(blank)

            bg = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_left, Inches(0), panel_w, full_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
            bg.line.fill.background()

            sld.shapes.add_picture(img_path, Inches(0), Inches(0), width=img_w, height=full_h)

            title_box = sld.shapes.add_textbox(Inches(5.95), Inches(0.35), Inches(3.9), Inches(0.85))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_meta.get("overlay_title", "")
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(30, 41, 59)
            try:
                p.font.name = "PingFang SC"
            except Exception:
                pass
            p.alignment = PP_ALIGN.LEFT

            body_box = sld.shapes.add_textbox(Inches(5.95), Inches(1.35), Inches(3.9), Inches(3.95))
            btf = body_box.text_frame
            btf.word_wrap = True
            btf.vertical_anchor = MSO_ANCHOR.TOP
            body_text = slide_meta.get("overlay_body", "")
            first_line = True
            for line in body_text.split("\n"):
                if first_line:
                    para = btf.paragraphs[0]
                    para.text = line
                    first_line = False
                else:
                    para = btf.add_paragraph()
                    para.text = line
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(51, 65, 85)
                para.space_after = Pt(6)
                try:
                    para.font.name = "PingFang SC"
                except Exception:
                    pass

        out_pptx = deck_dir / pptx_filename
        prs.save(str(out_pptx))
        logger.info("%s✓ 图文并茂 PPTX: %s (%.1f MB)", prefix, out_pptx.name, out_pptx.stat().st_size / (1024 * 1024))
        logger.info("")
    except Exception as e:
        logger.error("%sPPTX 失败: %s", prefix, e)
        import traceback
        traceback.print_exc()
        return False

    return True


def main() -> bool:
    parser = argparse.ArgumentParser(description="DeerFlow 本地 SDXL 生成 PPTX（单套 Tesla 或多套博世/自动驾驶主题）")
    parser.add_argument(
        "--mode",
        choices=("batch", "tesla", "illustrated", "logic"),
        default="batch",
        help="batch=4 套博世 + 归档；tesla=Tesla；illustrated=左图右文；logic=逻辑推理可视化 3 页全图",
    )
    args = parser.parse_args()

    logger.info("=" * 70)
    logger.info("DeerFlow PPT 生成 - 改进版本 (mode=%s)", args.mode)
    logger.info("=" * 70)
    logger.info("开始时间: %s", datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
    logger.info("")

    path_mapper = PathMapper()
    workspace = path_mapper.cache_dir / "workspace"
    outputs_root = path_mapper.cache_dir / "outputs"
    workspace.mkdir(parents=True, exist_ok=True)
    outputs_root.mkdir(parents=True, exist_ok=True)

    logger.info("工作目录: %s", workspace)
    logger.info("输出根目录: %s", outputs_root)
    logger.info("")

    script_dir = Path(__file__).resolve().parent

    batch_root = outputs_root / "bosch_automotive_batch"
    if args.mode == "batch":
        export_bosch_batch_12_prompts(workspace, script_dir)
        logger.info("")

    if args.mode == "batch":
        decks = BATCH_DECKS
        batch_root.mkdir(parents=True, exist_ok=True)
        logger.info("批量模式: %s 套演示，每套 3 页，输出: %s", len(decks), batch_root)
        logger.info("")
    elif args.mode == "tesla":
        decks = [{"id": "tesla_fsd", "pptx_name": "Tesla_FSD_Presentation.pptx", "plan": TESLA_LEGACY_PLAN}]
    else:
        decks = []

    # [1/4] 计划说明
    logger.info("=" * 70)
    logger.info("[1/4] 演示计划")
    logger.info("=" * 70)
    if args.mode == "batch":
        for d in BATCH_DECKS:
            logger.info("  - %s → %s", d["id"], d["pptx_name"])
    elif args.mode == "tesla":
        logger.info("  - Tesla FSD (legacy) → Tesla_FSD_Presentation.pptx")
    elif args.mode == "illustrated":
        logger.info("  - illustrated → %s (左图右文)", ILLUSTRATED_AD_DECK["pptx_name"])
    elif args.mode == "logic":
        logger.info("  - logic → %s (逻辑推理可视化)", LOGIC_REASONING_DECK["pptx_name"])
    logger.info("")

    # [2/4] 模型
    logger.info("=" * 70)
    logger.info("[2/4] 加载 SDXL 模型（仅一次）")
    logger.info("=" * 70)
    try:
        model_manager = ModelManager()
        model_manager.load_model()
    except Exception as e:
        logger.error("模型加载失败: %s", e)
        return False
    logger.info("")

    all_ok = True

    if args.mode == "illustrated":
        ill_dir = outputs_root / "ad_illustrated_zh"
        ok = run_one_deck_illustrated(ILLUSTRATED_AD_DECK, ill_dir, model_manager, deck_label="illustrated")
        all_ok = ok
    elif args.mode == "logic":
        ld = LOGIC_REASONING_DECK
        logic_dir = outputs_root / "logic_reasoning_viz"
        ok = run_one_deck(ld["plan"], logic_dir, ld["pptx_name"], model_manager, deck_label="logic")
        all_ok = ok
    else:
        for idx, deck in enumerate(decks, start=1):
            deck_id = deck["id"]
            ppt_plan = deck["plan"]
            pptx_name = deck["pptx_name"]
            if args.mode == "batch":
                deck_dir = batch_root / f"{idx:02d}_{_slug(deck_id)}"
                label = f"{idx}/{len(decks)} {deck_id}"
            else:
                deck_dir = outputs_root
                label = "tesla"

            logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>")
            logger.info("套 %s/%s: %s", idx, len(decks), ppt_plan.get("title", deck_id))
            logger.info(">>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>>")

            ok = run_one_deck(ppt_plan, deck_dir, pptx_name, model_manager, deck_label=label)
            if not ok:
                all_ok = False

    logger.info("=" * 70)
    if all_ok:
        logger.info("✅ 全部 PPT 任务完成 (mode=%s)", args.mode)
    else:
        logger.error("⚠️ 部分套失败，请向上查看日志")
    logger.info("=" * 70)
    if args.mode == "batch":
        logger.info("批量输出目录: %s", batch_root)
        logger.info("12 组 prompt 归档: %s / %s", workspace / "bosch_batch_12_prompts.json", script_dir / "bosch_batch_12_prompts.json")
    elif args.mode == "illustrated":
        logger.info("图文并茂输出: %s", outputs_root / "ad_illustrated_zh")
    elif args.mode == "logic":
        logger.info("逻辑推理可视化输出: %s", outputs_root / "logic_reasoning_viz")
    else:
        logger.info("输出目录: %s", outputs_root)
    logger.info("💰 本地 SDXL，无 API 计费")
    logger.info("")
    return all_ok


if __name__ == "__main__":
    try:
        success = main()
        sys.exit(0 if success else 1)
    except Exception as e:
        logger.error(f"致命错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
