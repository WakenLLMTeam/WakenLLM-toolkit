#!/usr/bin/env python3
"""
Tesla FSD PPT Generator - Direct Code Execution
直接在代码中调用 DeerFlow agent 生成 3 页 PPT，无需前端
"""

import sys
import asyncio
from pathlib import Path
from datetime import datetime

# 添加 DeerFlow 到路径
sys.path.insert(0, str(Path(__file__).parent / "backend" / "packages" / "harness"))

from deerflow.skills import load_skills


async def generate_fsd_ppt():
    """直接调用 skills 生成 PPT"""
    
    print("=" * 80)
    print("Tesla FSD PPT Generation - Direct Code Execution")
    print("=" * 80)
    print()
    
    # 准备工作目录
    workspace = Path("/tmp/fsd_ppt_direct/workspace")
    outputs = Path("/tmp/fsd_ppt_direct/outputs")
    workspace.mkdir(parents=True, exist_ok=True)
    outputs.mkdir(parents=True, exist_ok=True)
    
    print(f"[1/5] Setup")
    print(f"      Workspace: {workspace}")
    print(f"      Outputs:   {outputs}")
    print()
    
    # 加载 skills
    print("[2/5] Loading skills...")
    try:
        skills = load_skills(enabled_only=True)
        print(f"      ✓ Loaded {len(skills)} skills")
        
        ppt_skill = next((s for s in skills if s.name == 'ppt-generation'), None)
        img_skill = next((s for s in skills if s.name == 'image-generation'), None)
        
        if ppt_skill:
            print(f"      ✓ PPT Generation: {ppt_skill.name}")
        if img_skill:
            print(f"      ✓ Image Generation: {img_skill.name}")
    except Exception as e:
        print(f"      Error loading skills: {e}")
        return False
    
    print()
    
    # 创建 PPT 规划
    print("[3/5] Creating PPT plan...")
    ppt_plan = {
        "title": "Tesla Full Self-Driving (FSD)",
        "style": "gradient-modern",
        "aspect_ratio": "16:9",
        "slides": [
            {
                "slide_number": 1,
                "type": "title",
                "title": "Tesla Full Self-Driving",
                "subtitle": "Autonomous Driving Revolution",
                "visual_description": "Bold gradient background flowing from deep Tesla red (#E82127) through orange to electric blue (#0071E3). Center: large metallic silver 'FSD' text with glow effect. Behind: abstract neural network with glowing nodes and connection lines. Top right: Tesla Model 3 silhouette. Bottom: digital highway with autonomous path overlay. Professional, premium, futuristic."
            },
            {
                "slide_number": 2,
                "type": "content",
                "title": "Key Technologies",
                "key_points": [
                    "8 cameras + 12 ultrasonic sensors",
                    "Real-time neural network processing",
                    "Fleet Learning: 4+ billion miles",
                    "Safety-first redundancy architecture"
                ],
                "visual_description": "Tech specs slide. Center: Tesla Model 3 top-down view with sensor locations marked. Around: four floating glass cards showing camera array, processing speed, fleet learning stats, safety features. Gradient background, premium tech aesthetic. Data visualization elements."
            },
            {
                "slide_number": 3,
                "type": "conclusion",
                "title": "The Future of Mobility",
                "subtitle": "Join the autonomous revolution",
                "visual_description": "Dramatic finale: multiple Tesla vehicles on modern highway at sunset. Tech overlays and glowing elements. Bold 'The Future is Autonomous' text overlay. Tesla logo and stats: 8M+ vehicles, 4B+ miles, 99.9% safety rating. Inspirational cinematography, premium presentation quality."
            }
        ]
    }
    
    plan_file = workspace / "fsd_ppt_plan.json"
    import json
    with open(plan_file, "w") as f:
        json.dump(ppt_plan, f, indent=2)
    
    print(f"      ✓ Plan created: {plan_file}")
    print(f"      ✓ Slides: {len(ppt_plan['slides'])}")
    print()
    
    # 生成图像
    print("[4/5] Generating slide images with local SDXL...")
    print("      (This may take 5-15 minutes depending on hardware)")
    print()
    
    try:
        import torch
        from diffusers import StableDiffusionXLPipeline, DPMSolverMultistepScheduler
        
        device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"      Using device: {device}")
        
        # 加载模型
        print("      Loading SDXL model...")
        pipe = StableDiffusionXLPipeline.from_pretrained(
            "stabilityai/stable-diffusion-xl-base-1.0",
            torch_dtype=torch.float16 if device == "cuda" else torch.float32,
            use_safetensors=True,
            variant="fp16" if device == "cuda" else None,
        )
        
        pipe.scheduler = DPMSolverMultistepScheduler.from_config(pipe.scheduler.config)
        pipe = pipe.to(device)
        
        if device == "cpu":
            pipe.enable_sequential_cpu_offload()
        
        # 生成 3 张幻灯片
        num_steps = 15 if device == "cpu" else 20
        
        for slide in ppt_plan['slides']:
            slide_num = slide['slide_number']
            print(f"      Generating slide {slide_num}/3...")
            
            prompt = slide['visual_description']
            
            try:
                with torch.inference_mode():
                    image = pipe(
                        prompt=prompt,
                        height=576,
                        width=1024,
                        num_inference_steps=num_steps,
                        guidance_scale=7.5,
                    ).images[0]
                
                output_path = outputs / f"slide-{slide_num:02d}.jpg"
                image.save(output_path, quality=95)
                print(f"      ✓ Saved: slide-{slide_num:02d}.jpg")
            except Exception as e:
                print(f"      ✗ Error generating slide {slide_num}: {e}")
                return False
        
        del pipe
        if device == "cuda":
            torch.cuda.empty_cache()
        
    except ImportError:
        print("      ✗ torch/diffusers not available")
        print("      Note: Ensure you've run: cd backend && uv add diffusers transformers torch")
        return False
    except Exception as e:
        print(f"      ✗ Error: {e}")
        return False
    
    print()
    
    # 组合成 PPTX
    print("[5/5] Composing into PPTX...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9
        prs.title = ppt_plan['title']
        
        # 添加每一张幻灯片
        slide_images = sorted(outputs.glob("slide-*.jpg"))
        
        for i, img_path in enumerate(slide_images, 1):
            # 添加空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加图像（全尺寸）
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(
                str(img_path),
                left, top,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            print(f"      ✓ Added slide {i}")
        
        # 保存 PPTX
        output_pptx = outputs / "Tesla_FSD_3Slides.pptx"
        prs.save(str(output_pptx))
        
        print(f"      ✓ PPTX saved: {output_pptx}")
        print(f"      ✓ File size: {output_pptx.stat().st_size / (1024*1024):.2f} MB")
        
    except ImportError:
        print("      ✗ python-pptx not available")
        print("      Slides generated but not composed into PPTX")
        return False
    except Exception as e:
        print(f"      ✗ Error composing PPTX: {e}")
        return False
    
    print()
    print("=" * 80)
    print("✓ PPT Generation Complete!")
    print("=" * 80)
    print()
    print(f"Output file: {output_pptx}")
    print(f"  - 3 slides with AI-generated images")
    print(f"  - Gradient-modern style (red → blue)")
    print(f"  - 1024×576 resolution (16:9)")
    print(f"  - Ready to download and present")
    print()
    print(f"Timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    return True


def main():
    """Main entry point"""
    try:
        success = asyncio.run(generate_fsd_ppt())
        return 0 if success else 1
    except KeyboardInterrupt:
        print("\n\nInterrupted by user")
        return 1
    except Exception as e:
        print(f"\n✗ Error: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
