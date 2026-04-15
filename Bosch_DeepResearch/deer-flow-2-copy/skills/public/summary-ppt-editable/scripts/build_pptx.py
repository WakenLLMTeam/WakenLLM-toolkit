#!/usr/bin/env python3
"""
Build editable-text PPTX from a JSON plan (summary-ppt-editable skill).

Usage:
  python build_pptx.py --plan-file plan.json --output-file out.pptx
"""

from __future__ import annotations

import argparse
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import math as _math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ── Text-fit utility ──────────────────────────────────────────────────────────

def _fit_text_pt(
    text: str,
    box_w_in: float,
    box_h_in: float,
    start_pt: float,
    min_pt: float = 9.0,
    line_spacing: float = 1.35,
) -> int:
    """
    Return largest integer pt ≤ start_pt such that `text` fits in a word-wrapped
    box of box_w_in × box_h_in inches.

    CJK characters (U+2E80+) count as 1.8× the width of an ASCII character.
    Heuristic char width ≈ 0.55 × pt / 72 inches (proportional font, average).
    """
    if not text or box_w_in <= 0 or box_h_in <= 0:
        return int(max(start_pt, min_pt))
    pt = float(start_pt)
    while pt > min_pt:
        char_w = 0.55 * pt / 72
        visual_len = sum(1.8 if ord(c) > 0x2E80 else 1.0 for c in text)
        chars_per_line = max(1.0, (box_w_in * 0.90) / char_w)
        n_lines = _math.ceil(visual_len / chars_per_line)
        line_h = line_spacing * pt / 72
        if n_lines * line_h <= box_h_in * 0.92:
            break
        pt -= 0.5
    return int(max(pt, min_pt))


def _fit_bullets_pt(
    bullets: List[str],
    box_w_in: float,
    box_h_in: float,
    start_pt: float = 20.0,
    min_pt: float = 9.0,
    line_spacing: float = 1.4,
    space_after_ratio: float = 0.40,
) -> int:
    """
    Return largest integer pt such that all bullet lines fit in the box.
    Each bullet may word-wrap; space_after = pt * space_after_ratio / 72 inches.
    """
    lines = [(b or "").strip() for b in bullets if (b or "").strip()]
    if not lines or box_w_in <= 0 or box_h_in <= 0:
        return int(max(start_pt, min_pt))
    pt = float(start_pt)
    while pt > min_pt:
        char_w = 0.55 * pt / 72
        line_h = (line_spacing * pt + pt * space_after_ratio) / 72
        total_h = 0.0
        for text in lines:
            visual_len = sum(1.8 if ord(c) > 0x2E80 else 1.0 for c in text)
            chars_per_line = max(1.0, (box_w_in * 0.88) / char_w)
            n_wrap = _math.ceil(visual_len / chars_per_line)
            total_h += n_wrap * line_h
        if total_h <= box_h_in * 0.92:
            break
        pt -= 0.5
    return int(max(pt, min_pt))


def _add_picture_fit(slide, img_path: str, box_left, box_top, box_w, box_h):
    """
    Insert an image preserving its original aspect ratio, fitted (letterboxed)
    inside the given bounding box and centred within it.
    Avoids the distortion caused by passing both width and height to add_picture.

    Returns (actual_left, actual_top, actual_w, actual_h) — the real pixel
    bounds of the rendered image (after letterboxing), so callers can place
    captions or other elements tightly below/beside the image.
    """
    from PIL import Image as _PILImage
    with _PILImage.open(img_path) as im:
        img_w, img_h = im.size

    img_aspect = img_w / img_h
    box_aspect = box_w / box_h

    if img_aspect >= box_aspect:
        # Constrained by width
        fit_w = box_w
        fit_h = box_w / img_aspect
    else:
        # Constrained by height
        fit_h = box_h
        fit_w = box_h * img_aspect

    # Centre within box
    offset_left = (box_w - fit_w) / 2
    offset_top  = (box_h - fit_h) / 2

    actual_left = box_left + offset_left
    actual_top  = box_top  + offset_top

    slide.shapes.add_picture(
        img_path,
        actual_left,
        actual_top,
        width=fit_w,
        height=fit_h,
    )
    return actual_left, actual_top, fit_w, fit_h



def _truncate(text: str, max_chars: int) -> str:
    """Trim text to max_chars, appending '…' if truncated."""
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    # Try to cut at a word boundary
    cut = text[:max_chars - 1].rsplit(" ", 1)[0]
    return (cut if len(cut) > max_chars // 2 else text[:max_chars - 1]) + "…"


def _rgb(t: Optional[List[int]], default: Tuple[int, int, int]) -> RGBColor:
    if t and len(t) == 3:
        return RGBColor(int(t[0]), int(t[1]), int(t[2]))
    return RGBColor(default[0], default[1], default[2])


def _slide_size(aspect: str) -> Tuple[Any, Any]:
    if aspect == "4:3":
        return Inches(10), Inches(7.5)
    return Inches(13.333), Inches(7.5)


def _add_modules(
    text_frame,
    modules: List[Dict[str, Any]],
    *,
    font_name: str,
    body_rgb: RGBColor,
    heading_rgb: RGBColor,
    heading_pt: int,
    body_pt: int,
) -> None:
    """Multi-block text: each module has heading + bullets (editable)."""
    first = True
    for mod in modules:
        if not isinstance(mod, dict):
            continue
        h = (mod.get("heading") or "").strip()
        bs = mod.get("bullets") or []
        if h:
            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False
            p.text = _truncate(h, 25)
            p.font.bold = True
            p.font.size = Pt(heading_pt)
            p.font.color.rgb = heading_rgb
            p.space_after = Pt(4)
            try:
                p.font.name = font_name
            except Exception:
                pass
        for line in bs:
            line = (line or "").strip()
            if not line:
                continue
            p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            first = False
            p.text = _truncate(line, 50)
            p.level = 0
            p.font.size = Pt(body_pt)
            p.font.color.rgb = body_rgb
            p.space_after = Pt(6)
            try:
                p.font.name = font_name
            except Exception:
                pass


def _render_slide_body(text_frame, slide: Dict[str, Any], *, font_name: str, body_rgb: RGBColor, title_rgb: RGBColor) -> None:
    modules = slide.get("modules")
    if modules and isinstance(modules, list):
        _add_modules(
            text_frame,
            modules,
            font_name=font_name,
            body_rgb=body_rgb,
            heading_rgb=title_rgb,
            heading_pt=15,
            body_pt=14,
        )
    else:
        _add_bullets(text_frame, slide.get("bullets") or [], font_name=font_name, body_rgb=body_rgb, size_pt=16)


def _add_cards(
    sld,
    cards: List[Dict[str, Any]],
    *,
    left: Any,
    top: Any,
    total_w: Any,
    total_h: Any,
    font_name: str,
    body_rgb: RGBColor,
    heading_rgb: RGBColor,
    accent_rgb: RGBColor,
    card_bg_rgb: Optional[RGBColor] = None,
    cols: int = 0,
) -> None:
    """Render a list of cards as a grid of bordered text boxes.

    Each card dict supports:
      heading   – str, bold header shown at top of card (optional)
      bullets   – list[str], bullet lines (optional)
      icon      – str, single emoji / symbol prepended to heading (optional)
      bg_rgb    – list[int, int, int], per-card background override (optional)
      stat      – str, large highlight number/text shown in card footer (optional)
      stat_label– str, small label below stat (optional)

    Layout:
      - Card height is computed from actual content (heading + bullets) so
        cards never have excess empty space.
      - If the cards don't fill total_h, the leftover space below is used for
        a bottom accent strip with per-card stat values (if any card has stat).
      - cols=0 → auto: 1→1 col, 2→2, 3→3, 4→2×2, 5-6→3 cols, 7+→3 cols
    """
    n = len(cards)
    if n == 0:
        return

    if cols <= 0:
        if n <= 1:
            cols = 1
        elif n == 2:
            cols = 2
        elif n <= 3:
            cols = 3
        elif n == 4:
            cols = 2
        else:
            cols = 3
    rows = (n + cols - 1) // cols

    gap = Inches(0.12)
    card_w = (total_w - gap * (cols - 1)) / cols

    # ── Compute natural card height from content ──────────────────────────────
    # heading line ≈ Pt(13) * 1.4/72 in + space_after Pt(5)/72 ≈ 0.322in
    # bullet line  ≈ Pt(11) * 1.4/72 in + space_after Pt(4)/72  ≈ 0.269in
    # top bar = 0.055in, pad_y = 0.18in (top+bottom = 0.36in)
    HEADING_LINE_H = Inches(0.33)
    BULLET_LINE_H  = Inches(0.27)
    BAR_H          = Inches(0.055)
    PAD_Y          = Inches(0.18)
    PAD_X          = Inches(0.16)

    max_lines_per_row: List[int] = []
    for row_i in range(rows):
        max_lines = 0
        for col_i in range(cols):
            idx = row_i * cols + col_i
            if idx >= n:
                continue
            card = cards[idx]
            has_heading = bool((card.get("heading") or "").strip())
            n_bullets = len([b for b in (card.get("bullets") or []) if (b or "").strip()])
            lines = (1 if has_heading else 0) + n_bullets
            max_lines = max(max_lines, lines)
        max_lines_per_row.append(max_lines)

    def _natural_card_h(n_lines: int) -> Any:
        has_h = n_lines > 0
        bullet_lines = max(0, n_lines - 1) if has_h else n_lines
        return BAR_H + PAD_Y + HEADING_LINE_H + bullet_lines * BULLET_LINE_H + PAD_Y

    row_heights = [_natural_card_h(m) for m in max_lines_per_row]

    # Cap: don't exceed the slot height / rows (allow slight overflow into stat area)
    max_row_h = (total_h - gap * (rows - 1)) / rows
    row_heights = [min(h, max_row_h) for h in row_heights]

    total_card_h = sum(row_heights) + gap * (rows - 1)
    leftover = total_h - total_card_h  # space below cards

    # ── Draw cards ────────────────────────────────────────────────────────────
    default_bg = card_bg_rgb or RGBColor(245, 247, 250)
    row_tops: List[Any] = []
    cur_top = top
    for rh in row_heights:
        row_tops.append(cur_top)
        cur_top += rh + gap

    for idx, card in enumerate(cards):
        col_i = idx % cols
        row_i = idx // cols
        cx = left + col_i * (card_w + gap)
        cy = row_tops[row_i]
        ch = row_heights[row_i]

        bg = card.get("bg_rgb")
        bg_color = RGBColor(int(bg[0]), int(bg[1]), int(bg[2])) if bg and len(bg) == 3 else default_bg

        rect = sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, card_w, ch)
        rect.fill.solid()
        rect.fill.fore_color.rgb = bg_color
        rect.line.color.rgb = accent_rgb
        rect.line.width = Pt(1.0)
        rect.adjustments[0] = 0.05

        bar = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, card_w, BAR_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_rgb
        bar.line.fill.background()

        tb = sld.shapes.add_textbox(cx + PAD_X, cy + BAR_H + PAD_Y,
                                    card_w - PAD_X * 2, ch - BAR_H - PAD_Y * 2)
        tf = tb.text_frame
        tf.word_wrap = True

        bullets = [b for b in (card.get("bullets") or []) if (b or "").strip()]
        first = True
        heading = (card.get("heading") or "").strip()
        icon = (card.get("icon") or "").strip()
        if heading:
            label = f"{icon}  {_truncate(heading, 20)}" if icon else _truncate(heading, 20)
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = label
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = heading_rgb
            p.space_after = Pt(5)
            try:
                p.font.name = font_name
            except Exception:
                pass

        for line in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"• {_truncate(line, 45)}"
            p.font.size = Pt(11)
            p.font.color.rgb = body_rgb
            p.space_after = Pt(4)
            try:
                p.font.name = font_name
            except Exception:
                pass

    # ── Bottom stat strip (if leftover space ≥ 0.5in and any card has stat) ───
    has_stats = any((card.get("stat") or "").strip() for card in cards)
    if leftover >= Inches(0.45) and has_stats:
        strip_top = top + total_card_h + Inches(0.08)
        strip_h = leftover - Inches(0.08)
        # Draw a light accent background strip
        strip = sld.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, strip_top, total_w, strip_h,
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = RGBColor(
            accent_rgb.red, accent_rgb.green, accent_rgb.blue
        ) if hasattr(accent_rgb, 'red') else accent_rgb
        # Use very light tint
        strip.fill.fore_color.rgb = RGBColor(
            min(255, accent_rgb[0] + 220) if isinstance(accent_rgb, (list, tuple)) else 240,
            min(255, accent_rgb[1] + 220) if isinstance(accent_rgb, (list, tuple)) else 245,
            min(255, accent_rgb[2] + 220) if isinstance(accent_rgb, (list, tuple)) else 250,
        )
        strip.line.fill.background()

        stat_w = total_w / max(cols, 1)
        for col_i in range(cols):
            # find the last row's card for this column
            card_idx = (rows - 1) * cols + col_i
            if card_idx >= n:
                card_idx = col_i  # fallback to first row
            if card_idx >= n:
                continue
            card = cards[card_idx]
            stat = (card.get("stat") or "").strip()
            stat_label = (card.get("stat_label") or "").strip()
            if not stat:
                continue
            sx = left + col_i * stat_w
            stat_box = sld.shapes.add_textbox(sx + Inches(0.1), strip_top + Inches(0.04),
                                              stat_w - Inches(0.2), strip_h - Inches(0.08))
            stf = stat_box.text_frame
            stf.word_wrap = False
            sp = stf.paragraphs[0]
            sp.text = stat
            sp.font.bold = True
            sp.font.size = Pt(min(22, max(14, int(strip_h / 914400 * 48))))
            sp.font.color.rgb = heading_rgb
            sp.alignment = PP_ALIGN.CENTER
            try:
                sp.font.name = font_name
            except Exception:
                pass
            if stat_label:
                sp2 = stf.add_paragraph()
                sp2.text = stat_label
                sp2.font.size = Pt(9)
                sp2.font.color.rgb = body_rgb
                sp2.alignment = PP_ALIGN.CENTER
                try:
                    sp2.font.name = font_name
                except Exception:
                    pass


def _add_bullets(text_frame, bullets: List[str], *, font_name: str, body_rgb: RGBColor, size_pt: int) -> None:
    lines = [((x or "").strip()) for x in bullets if (x or "").strip()]
    if not lines:
        text_frame.paragraphs[0].text = ""
        return
    # space_after scales with font: 8pt at ≥18pt, 5pt at 13pt, 2pt at ≤10pt
    space_after = Pt(max(2, int(size_pt * 0.45)))
    first = True
    for line in lines:
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        p.text = "· " + _truncate(line, 58)
        p.level = 0
        p.font.size = Pt(size_pt)
        p.font.color.rgb = body_rgb
        p.space_after = space_after
        try:
            p.font.name = font_name
        except Exception:
            pass


# ──────────────────────────────────────────────────────────────────────────────
# Visual design helpers (styled text boxes)
# ──────────────────────────────────────────────────────────────────────────────

def _lighten_rgb(rgb: RGBColor, factor: float = 0.82) -> RGBColor:
    """Blend rgb toward white. factor=1.0 → pure white, 0.0 → unchanged.
    Works with RGBColor (tuple subclass) using index access."""
    r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    return RGBColor(
        min(255, int(r + (255 - r) * factor)),
        min(255, int(g + (255 - g) * factor)),
        min(255, int(b + (255 - b) * factor)),
    )


def _add_body_region(sld, left, top, width, height, accent_rgb: RGBColor):
    """Styled text region: left accent bar + surface background + thin border.
    Returns a text_frame (word_wrap=True, TOP anchor, inner margin set)."""
    BAR_W   = Inches(0.042)
    SURFACE = RGBColor(245, 247, 250)
    BORDER  = RGBColor(218, 224, 234)

    # Thin left accent bar
    bar = sld.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top + Inches(0.06), BAR_W, height - Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.fill.background()

    # Text box with surface fill (acts as background + border)
    tb = sld.shapes.add_textbox(
        left + BAR_W + Inches(0.01), top,
        width - BAR_W - Inches(0.01), height,
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = SURFACE
    tb.line.color.rgb = BORDER
    tb.line.width = Pt(0.75)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left   = Inches(0.14)
    tf.margin_right  = Inches(0.10)
    tf.margin_top    = Inches(0.10)
    tf.margin_bottom = Inches(0.08)
    return tf


def _add_module_cards(
    sld,
    left, top, width, height,
    modules: List[Dict[str, Any]],
    *,
    font_name: str,
    body_rgb: RGBColor,
    heading_rgb: RGBColor,
    accent_rgb: RGBColor,
    heading_pt: int = 13,
    body_pt: int = 12,
) -> None:
    """Render modules as horizontal rows: [pill badge LEFT] | [bullet text RIGHT].

    Each row:
      [accent bar] [heading badge centred vertically] | [divider] [bullet lines]
    """
    valid = [m for m in modules if isinstance(m, dict)]
    if not valid:
        return
    n = len(valid)
    SURFACE   = RGBColor(245, 247, 250)
    BORDER    = RGBColor(218, 224, 234)
    accent_lt = _lighten_rgb(accent_rgb, 0.82)

    gap   = Inches(0.08)
    row_h = (height - gap * (n - 1)) / max(n, 1)

    BAR_W     = Inches(0.038)
    LEFT_W    = int(width * 0.28)       # badge column width
    DIV_W     = Inches(0.007)
    RIGHT_OFF = LEFT_W + Inches(0.10)   # right text column start offset from `left`
    RIGHT_W   = width - RIGHT_OFF - Inches(0.05)

    for i, mod in enumerate(valid):
        rx = left
        ry = top + i * (row_h + gap)
        rh = row_h

        h_text  = (mod.get("heading") or "").strip()
        bullets = [b.strip() for b in (mod.get("bullets") or []) if (b or "").strip()]

        # ── Row background ────────────────────────────────────────────────────
        row_bg = sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, width, rh)
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = SURFACE
        row_bg.line.color.rgb = BORDER
        row_bg.line.width = Pt(0.75)
        row_bg.adjustments[0] = 0.04

        # Left accent bar
        bar = sld.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            rx + Inches(0.01), ry + Inches(0.05),
            BAR_W, rh - Inches(0.10),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_rgb
        bar.line.fill.background()

        # ── Heading pill badge (left column, vertically centred) ──────────────
        if h_text:
            BADGE_H   = Inches(0.30)
            badge_top = ry + (rh - BADGE_H) / 2
            badge_lx  = rx + BAR_W + Inches(0.06)
            badge_w   = LEFT_W - BAR_W - Inches(0.10)

            badge = sld.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                badge_lx, badge_top, badge_w, BADGE_H,
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_lt
            badge.line.fill.background()
            badge.adjustments[0] = 0.50  # full pill

            _badge_inner_w = max(0.3, float(badge_w - Inches(0.12)) / 914400)
            _badge_inner_h = max(0.15, float(BADGE_H) / 914400)
            _h_text_trunc = _truncate(h_text, 18)
            _badge_pt = _fit_text_pt(_h_text_trunc, _badge_inner_w, _badge_inner_h,
                                     start_pt=heading_pt, min_pt=9)
            htb = sld.shapes.add_textbox(
                badge_lx + Inches(0.06), badge_top,
                badge_w - Inches(0.06), BADGE_H,
            )
            htf = htb.text_frame
            htf.vertical_anchor = MSO_ANCHOR.MIDDLE
            hp  = htf.paragraphs[0]
            hp.text = _h_text_trunc
            hp.font.bold = True
            hp.font.size = Pt(_badge_pt)
            hp.font.color.rgb = heading_rgb
            hp.alignment = PP_ALIGN.CENTER
            try:
                hp.font.name = font_name
            except Exception:
                pass

        # ── Thin vertical divider ─────────────────────────────────────────────
        div = sld.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            rx + LEFT_W, ry + Inches(0.08),
            DIV_W, rh - Inches(0.16),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = BORDER
        div.line.fill.background()

        # ── Bullet text (right column) ────────────────────────────────────────
        txt_top = ry + Inches(0.09)
        txt_h   = rh - Inches(0.18)
        if txt_h > Inches(0.12):
            n_b = max(len(bullets), 1)
            natural_pt = float(txt_h) / 914400 * 0.78 * 72 / (n_b * 1.45)
            size_pt = int(max(10, min(body_pt, natural_pt)))

            btb = sld.shapes.add_textbox(rx + RIGHT_OFF, txt_top, RIGHT_W, txt_h)
            btf = btb.text_frame
            btf.word_wrap = True
            btf.vertical_anchor = MSO_ANCHOR.TOP
            first = True
            for line in bullets:
                p = btf.paragraphs[0] if first else btf.add_paragraph()
                first = False
                p.text = "· " + _truncate(line, 53)
                p.font.size = Pt(size_pt)
                p.font.color.rgb = body_rgb
                p.space_after = Pt(3)
                try:
                    p.font.name = font_name
                except Exception:
                    pass


def _render_content_body(
    sld,
    left, top, width, height,
    slide: Dict[str, Any],
    *,
    font_name: str,
    body_rgb: RGBColor,
    title_rgb: RGBColor,
    accent_rgb: RGBColor,
) -> None:
    """Dispatch slide body to module-cards or styled bullet region."""
    modules = slide.get("modules")
    if modules and isinstance(modules, list):
        _add_module_cards(
            sld, left, top, width, height, modules,
            font_name=font_name, body_rgb=body_rgb,
            heading_rgb=title_rgb, accent_rgb=accent_rgb,
        )
        return

    bullets = [(b or "").strip() for b in (slide.get("bullets") or []) if (b or "").strip()]
    box_w_in = float(width) / 914400
    box_h_in = float(height) / 914400
    # Inner text area shrinks by margins: ~0.14+0.10 = 0.24in horiz, ~0.10+0.08 = 0.18in vert
    inner_w = max(0.5, box_w_in - 0.24)
    inner_h = max(0.3, box_h_in - 0.18)
    size_pt = _fit_bullets_pt(bullets, inner_w, inner_h, start_pt=20, min_pt=9)
    tf = _add_body_region(sld, left, top, width, height, accent_rgb)
    _add_bullets(tf, bullets, font_name=font_name, body_rgb=body_rgb, size_pt=size_pt)


def build_pptx(plan: Dict[str, Any], output_file: str) -> str:
    aspect = plan.get("aspect_ratio", "16:9")
    slide_w, slide_h = _slide_size(aspect)
    theme = plan.get("theme") or {}
    accent = _rgb(theme.get("accent_rgb"), (0, 71, 227))
    body_rgb = _rgb(theme.get("body_rgb"), (51, 65, 85))
    title_rgb = _rgb(theme.get("title_rgb"), (15, 23, 42))

    font_title = theme.get("font_title", "PingFang SC")
    font_body = theme.get("font_body", "PingFang SC")

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    try:
        prs.core_properties.title = plan.get("title", "Presentation")
        prs.core_properties.author = theme.get("author", "DeerFlow summary-ppt-editable")
    except Exception:
        pass

    blank = prs.slide_layouts[6]
    slides_in = sorted(plan.get("slides", []), key=lambda s: int(s.get("slide_number", 0)))

    _section_counter = 0  # counts section slides only

    for slide in slides_in:
        stype = (slide.get("type") or "content").lower()
        if stype == "section":
            _section_counter += 1
        sld = prs.slides.add_slide(blank)

        if stype == "title":
            # ── Bottom accent band ────────────────────────────────────────────
            band_h = Inches(1.55)
            band = sld.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), slide_h - band_h, slide_w, band_h,
            )
            band.fill.solid()
            band.fill.fore_color.rgb = accent
            band.line.fill.background()

            # Small accent dot (top-left decoration)
            dot = sld.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.45), Inches(0.38), Inches(0.14), Inches(0.14),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = accent
            dot.line.fill.background()

            # Title (upper 60% of slide)
            _title_text = _truncate(slide.get("title", ""), 40)
            _title_box_w = float(slide_w - Inches(2.0)) / 914400
            _title_box_h = 2.2
            _title_pt = _fit_text_pt(_title_text, _title_box_w, _title_box_h, start_pt=40, min_pt=20)
            box = sld.shapes.add_textbox(Inches(1.0), Inches(1.5), slide_w - Inches(2.0), Inches(2.2))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = _title_text
            p.font.bold = True
            p.font.size = Pt(_title_pt)
            p.font.color.rgb = title_rgb
            p.alignment = PP_ALIGN.CENTER
            try:
                p.font.name = font_title
            except Exception:
                pass

            # Central accent divider line below title
            div = sld.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(slide_w * 0.30), Inches(3.72), int(slide_w * 0.40), Inches(0.04),
            )
            div.fill.solid()
            div.fill.fore_color.rgb = accent
            div.line.fill.background()

            # Subtitle inside accent band (white text)
            sub = slide.get("subtitle") or ""
            if sub:
                _sub_text = _truncate(sub, 70)
                _sub_pt = _fit_text_pt(_sub_text,
                                       float(slide_w - Inches(2.0)) / 914400, 0.88,
                                       start_pt=18, min_pt=11)
                sub_box = sld.shapes.add_textbox(
                    Inches(1.0), slide_h - band_h + Inches(0.38),
                    slide_w - Inches(2.0), Inches(0.88),
                )
                stf = sub_box.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                sp.text = _sub_text
                sp.font.size = Pt(_sub_pt)
                sp.font.color.rgb = RGBColor(255, 255, 255)
                sp.alignment = PP_ALIGN.CENTER
                try:
                    sp.font.name = font_body
                except Exception:
                    pass
        elif stype == "section":
            # ── Full-width top accent band ────────────────────────────────────
            top_band = sld.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_w, Inches(0.13),
            )
            top_band.fill.solid()
            top_band.fill.fore_color.rgb = accent
            top_band.line.fill.background()

            # Large faded section counter (decorative background)
            # Uses section order (01, 02, 03…) not slide page number
            sec_num = _section_counter
            if sec_num:
                num_str = str(sec_num).zfill(2)
                nb = sld.shapes.add_textbox(
                    slide_w - Inches(5.2), Inches(0.5), Inches(5.0), Inches(5.8),
                )
                nbtf = nb.text_frame
                nbp  = nbtf.paragraphs[0]
                nbp.text = num_str
                nbp.font.bold = True
                nbp.font.size = Pt(220)
                nbp.font.color.rgb = RGBColor(238, 241, 248)
                nbp.alignment = PP_ALIGN.RIGHT
                try:
                    nbp.font.name = font_title
                except Exception:
                    pass

            # Section title (left-aligned, center-vertical)
            _sec_title = _truncate(slide.get("title", ""), 40)
            _sec_box_w = float(int(slide_w * 0.68)) / 914400
            _sec_pt = _fit_text_pt(_sec_title, _sec_box_w, 1.6, start_pt=34, min_pt=18)
            box = sld.shapes.add_textbox(
                Inches(0.85), Inches(2.55), int(slide_w * 0.68), Inches(1.6),
            )
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = _sec_title
            p.font.bold = True
            p.font.size = Pt(_sec_pt)
            p.font.color.rgb = title_rgb
            try:
                p.font.name = font_title
            except Exception:
                pass

            # Short accent underline below section title
            uline = sld.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.85), Inches(4.18), Inches(1.6), Inches(0.055),
            )
            uline.fill.solid()
            uline.fill.fore_color.rgb = accent
            uline.line.fill.background()
        else:
            # content | summary | cards
            title_h = Inches(0.95)
            _ct_text = _truncate(slide.get("title", ""), 40)
            _ct_box_w = float(slide_w - Inches(1.1)) / 914400
            _ct_pt = _fit_text_pt(_ct_text, _ct_box_w, 0.95, start_pt=26, min_pt=14)
            tb = sld.shapes.add_textbox(Inches(0.55), Inches(0.35), slide_w - Inches(1.1), title_h)
            tfp = tb.text_frame
            tfp.word_wrap = True
            pp = tfp.paragraphs[0]
            pp.text = _ct_text
            pp.font.bold = True
            pp.font.size = Pt(_ct_pt)
            pp.font.color.rgb = title_rgb
            try:
                pp.font.name = font_title
            except Exception:
                pass

            # (accent underline below slide title removed)

            # ── Cards layout ──────────────────────────────────────────────
            cards = slide.get("cards")
            if cards and isinstance(cards, list):
                left_margin = Inches(0.55)
                top_body = Inches(1.35)
                _add_cards(
                    sld,
                    cards,
                    left=left_margin,
                    top=top_body,
                    total_w=slide_w - Inches(1.1),
                    total_h=slide_h - top_body - Inches(0.25),
                    font_name=font_body,
                    body_rgb=body_rgb,
                    heading_rgb=title_rgb,
                    accent_rgb=accent,
                    cols=int(slide.get("cards_cols", 0)),
                )
                # Notes
                notes = slide.get("notes")
                if notes:
                    try:
                        ns = sld.notes_slide
                        ns.notes_text_frame.text = str(notes)[:300]
                    except Exception:
                        pass
                continue  # skip figure logic below for cards slides

            fig = slide.get("figure") or {}
            img_path = fig.get("image_path") or ""
            has_figure = bool(img_path) or bool(fig.get("caption"))
            pos = (fig.get("position") or "right").lower()
            # Comparison tables and flowcharts always get full-slide layout
            viz_type = (fig.get("viz") or {}).get("type", "")
            if viz_type in ("comparison", "flowchart", "arch"):
                pos = "full"
            elif viz_type in ("timeline", "pipeline"):
                pos = "bottom"
            elif viz_type == "radar":
                pos = "right"
            if pos not in ("right", "bottom", "full"):
                pos = "right"

            left_margin = Inches(0.55)
            top_body = Inches(1.35)
            full_text_w = slide_w - Inches(1.1)

            # ── New layout dispatch ───────────────────────────────────────────
            layout = (slide.get("layout") or "default").lower()

            # Adaptive bullet font: use _fit_bullets_pt for accurate CJK-aware sizing
            def _bullet_pt(box_h_in: float, box_w_in: float = 0) -> int:
                buls = [b for b in (slide.get("bullets") or []) if (b or "").strip()]
                _w = box_w_in or float(full_text_w) / 914400
                return _fit_bullets_pt(buls, _w, box_h_in, start_pt=20, min_pt=9)

            def _caption_box(sld, actual, fig, font_body, body_rgb):
                cap = fig.get("caption")
                if not cap or not actual:
                    return
                _al, _at, _aw, _ah = actual
                cb = sld.shapes.add_textbox(_al, _at + _ah + Inches(0.05), _aw, Inches(0.38))
                ctf = cb.text_frame
                cp = ctf.paragraphs[0]
                cp.text = _truncate(str(cap), 60)
                cp.font.size = Pt(9)
                cp.font.color.rgb = body_rgb
                cp.alignment = PP_ALIGN.CENTER
                try:
                    cp.font.name = font_body
                except Exception:
                    pass

            # ── two_col: dual-column text, no viz ─────────────────────────────
            if layout == "two_col":
                all_bullets = [b for b in (slide.get("bullets") or []) if (b or "").strip()]
                mid = (len(all_bullets) + 1) // 2
                col_w = (full_text_w - Inches(0.45)) / 2
                body_h = slide_h - top_body - Inches(0.3)
                col_w_in = float(col_w) / 914400
                bpt = _bullet_pt(float(body_h) / 914400, col_w_in)
                if all_bullets[:mid]:
                    tf1 = _add_body_region(sld, left_margin, top_body, col_w, body_h, accent)
                    _add_bullets(tf1, all_bullets[:mid], font_name=font_body, body_rgb=body_rgb, size_pt=bpt)
                # accent divider
                div_x = left_margin + col_w + Inches(0.20)
                div = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, div_x,
                                           top_body + Inches(0.1), Inches(0.018), body_h - Inches(0.2))
                div.fill.solid()
                div.fill.fore_color.rgb = accent
                div.line.fill.background()
                if all_bullets[mid:]:
                    tf2 = _add_body_region(sld, div_x + Inches(0.25), top_body, col_w, body_h, accent)
                    _add_bullets(tf2, all_bullets[mid:], font_name=font_body, body_rgb=body_rgb, size_pt=bpt)

            # ── hero: viz dominant right (~65%), compact text left (~35%) ─────
            elif layout == "hero" and has_figure and img_path and os.path.isfile(img_path):
                txt_w = Inches(4.2)
                viz_x = left_margin + txt_w + Inches(0.25)
                viz_w = slide_w - viz_x - Inches(0.25)
                viz_top = Inches(1.15)
                viz_h = slide_h - viz_top - Inches(0.28)
                body_h = slide_h - top_body - Inches(0.3)
                _render_content_body(sld, left_margin, top_body, txt_w, body_h, slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)
                _actual = _add_picture_fit(sld, img_path, viz_x, viz_top, viz_w, viz_h)
                _caption_box(sld, _actual, fig, font_body, body_rgb)

            # ── split: balanced 50/50 ─────────────────────────────────────────
            elif layout == "split" and has_figure and img_path and os.path.isfile(img_path):
                txt_w = slide_w * 0.47 - left_margin
                viz_x = left_margin + txt_w + Inches(0.25)
                viz_w = slide_w - viz_x - Inches(0.25)
                body_h = slide_h - top_body - Inches(0.3)
                _render_content_body(sld, left_margin, top_body, txt_w, body_h, slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)
                _actual = _add_picture_fit(sld, img_path, viz_x, Inches(1.15), viz_w, slide_h - Inches(1.45))
                _caption_box(sld, _actual, fig, font_body, body_rgb)

            # ── viz_left: viz on left (~45%), text on right ───────────────────
            elif layout == "viz_left" and has_figure and img_path and os.path.isfile(img_path):
                viz_w = Inches(5.8)
                txt_x = left_margin + viz_w + Inches(0.25)
                txt_w = slide_w - txt_x - Inches(0.25)
                body_h = slide_h - top_body - Inches(0.3)
                _actual = _add_picture_fit(sld, img_path, left_margin, Inches(1.15), viz_w, slide_h - Inches(1.45))
                _caption_box(sld, _actual, fig, font_body, body_rgb)
                _render_content_body(sld, txt_x, top_body, txt_w, body_h, slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)

            # ── default / fallback: existing position-based dispatch ──────────
            # ----- Full-slide: figure spans the entire content area (comparison/flowchart) -----
            elif has_figure and pos == "full":
                # Optional compact summary text above the figure (1 line / short bullets)
                summary_lines = []
                if slide.get("bullets"):
                    summary_lines = [(b or "").strip() for b in slide["bullets"] if (b or "").strip()]
                elif slide.get("modules"):
                    for mod in slide["modules"]:
                        h = (mod.get("heading") or "").strip()
                        if h:
                            summary_lines.append(h)
                has_summary = bool(summary_lines)
                if has_summary:
                    sum_h = Inches(0.55)
                    sum_box = sld.shapes.add_textbox(left_margin, Inches(1.10), full_text_w, sum_h)
                    sf = sum_box.text_frame
                    sf.word_wrap = True
                    first_s = True
                    for line in summary_lines[:3]:  # at most 3 lines to keep compact
                        sp = sf.paragraphs[0] if first_s else sf.add_paragraph()
                        first_s = False
                        sp.text = line
                        sp.font.size = Pt(11)
                        sp.font.color.rgb = body_rgb
                        try:
                            sp.font.name = font_body
                        except Exception:
                            pass
                    fig_top = Inches(1.70)
                else:
                    fig_top = Inches(1.20)
                fig_h = slide_h - fig_top - Inches(0.35)
                fig_left = left_margin
                fig_w = full_text_w
                _actual = None
                if img_path and os.path.isfile(img_path):
                    _actual = _add_picture_fit(sld, img_path, fig_left, fig_top, fig_w, fig_h)
                cap = fig.get("caption")
                if cap:
                    if _actual:
                        _al, _at, _aw, _ah = _actual
                        cap_top = _at + _ah + Inches(0.05)
                    else:
                        cap_top = slide_h - Inches(0.38)
                        _al, _aw = fig_left, fig_w
                    cb = sld.shapes.add_textbox(_al, cap_top, _aw, Inches(0.32))
                    ctf = cb.text_frame
                    cp = ctf.paragraphs[0]
                    cp.text = str(cap)
                    cp.font.size = Pt(9)
                    cp.font.color.rgb = body_rgb
                    cp.alignment = PP_ALIGN.CENTER

            # ----- Bottom band: timeline / pipeline (text above, figure full width below) -----
            elif has_figure and pos == "bottom":
                # Count actual text lines to compute body height
                _content_lines = 0
                if slide.get("modules"):
                    for _m in slide["modules"]:
                        if (_m.get("heading") or "").strip():
                            _content_lines += 1
                        _content_lines += len([b for b in (_m.get("bullets") or []) if (b or "").strip()])
                else:
                    _content_lines = len([b for b in (slide.get("bullets") or []) if (b or "").strip()])

                cap = fig.get("caption")
                cap_h = Inches(0.35) if cap else Inches(0)
                gap = Inches(0.15)

                # Total available height from title bottom to slide bottom
                avail_h = slide_h - top_body - cap_h - Inches(0.08)

                if _content_lines == 0:
                    # No text: give everything to the figure
                    body_h = Inches(0)
                    fig_top = top_body
                    fig_h = avail_h
                else:
                    avail_in = float(avail_h) / 914400

                    if slide.get("modules") and isinstance(slide.get("modules"), list):
                        # Horizontal row layout: each row height = max(badge_h, bullets_h) + pad
                        # badge = 0.30in, bullets ≈ n*0.26in, row pad = 0.18in
                        valid_m = [m for m in slide["modules"] if isinstance(m, dict)]
                        n_m = max(len(valid_m), 1)
                        max_b = max(
                            (len([b for b in (m.get("bullets") or []) if (b or "").strip()]) for m in valid_m),
                            default=0,
                        )
                        row_min_h = max(0.30, max(1, max_b) * 0.26) + 0.18
                        natural_body = n_m * row_min_h + (n_m - 1) * 0.08
                        max_body_in = avail_in * 0.50
                    else:
                        natural_body = 0.25 + _content_lines * 0.26
                        max_body_in = avail_in * 0.35

                    body_h = Inches(max(0.55, min(max_body_in, natural_body)))
                    fig_top = top_body + body_h + gap
                    fig_h = slide_h - fig_top - cap_h - Inches(0.08)
                    fig_h = max(fig_h, Inches(1.5))

                # If image exists, respect its aspect ratio to avoid squashing
                if img_path and os.path.isfile(img_path):
                    try:
                        from PIL import Image as _PIL
                        with _PIL.open(img_path) as _im:
                            _iw, _ih = _im.size
                        img_ratio = _iw / _ih  # e.g. timeline ~3.5:1
                        slot_w = float(full_text_w) / 914400
                        # Ideal height for this image at full slot width
                        ideal_h_in = slot_w / img_ratio
                        ideal_h = Inches(ideal_h_in)
                        # If ideal height is less than what we allocated, shrink fig_h
                        # and push fig_top down (more text room or white space)
                        if ideal_h < fig_h * 0.92:
                            # re-center: push fig_top down a bit
                            freed = fig_h - ideal_h
                            fig_top = fig_top + freed * 0.3  # partial centering
                            fig_h = ideal_h
                    except Exception:
                        pass

                if body_h > Inches(0.05):
                    _render_content_body(sld, left_margin, top_body, full_text_w, body_h, slide,
                                         font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)

                pic_w = full_text_w
                pic_left = left_margin
                _actual_b = None
                if img_path and os.path.isfile(img_path):
                    _actual_b = _add_picture_fit(sld, img_path, pic_left, fig_top, pic_w, fig_h)
                else:
                    ph = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, pic_left, fig_top, pic_w, fig_h)
                    ph.fill.background()
                    ph.line.color.rgb = accent
                    ph.line.width = Pt(1.5)
                    tbox = sld.shapes.add_textbox(pic_left, fig_top + fig_h * 0.38, pic_w, Inches(0.85))
                    tfph = tbox.text_frame
                    tfp0 = tfph.paragraphs[0]
                    tfp0.text = fig.get("placeholder_text") or "路线 / 时间线配图区"
                    tfp0.font.size = Pt(11)
                    tfp0.font.color.rgb = body_rgb
                    tfp0.alignment = PP_ALIGN.CENTER
                if cap:
                    if _actual_b:
                        _al, _at, _aw, _ah = _actual_b
                        cap_top = _at + _ah + Inches(0.05)
                    else:
                        cap_top = fig_top + fig_h + Inches(0.04)
                        _al, _aw = pic_left, pic_w
                    cb = sld.shapes.add_textbox(_al, cap_top, _aw, cap_h)
                    ctf = cb.text_frame
                    cp = ctf.paragraphs[0]
                    cp.text = _truncate(str(cap), 60)
                    cp.font.size = Pt(10)
                    cp.font.color.rgb = body_rgb
                    cp.alignment = PP_ALIGN.CENTER
            elif has_figure and img_path and os.path.isfile(img_path):
                text_w = slide_w - Inches(6.1)
                pic_left = slide_w - Inches(5.35)
                pic_top = Inches(1.25)
                pic_w = Inches(4.75)
                pic_h = Inches(4.85)
                _render_content_body(sld, left_margin, top_body, text_w, slide_h - Inches(1.65), slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)
                _actual_r = _add_picture_fit(sld, img_path, pic_left, pic_top, pic_w, pic_h)
                cap = fig.get("caption")
                if cap:
                    _al, _at, _aw, _ah = _actual_r
                    cb = sld.shapes.add_textbox(_al, _at + _ah + Inches(0.05), _aw, Inches(0.45))
                    ctf = cb.text_frame
                    cp = ctf.paragraphs[0]
                    cp.text = str(cap)
                    cp.font.size = Pt(10)
                    cp.font.color.rgb = body_rgb
                    cp.alignment = PP_ALIGN.CENTER
            elif has_figure:
                text_w = slide_w - Inches(6.1)
                ph_left = slide_w - Inches(5.35)
                ph_top = Inches(1.25)
                ph_w = Inches(4.75)
                ph_h = Inches(4.85)
                _render_content_body(sld, left_margin, top_body, text_w, slide_h - Inches(1.65), slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)
                ph = sld.shapes.add_shape(MSO_SHAPE.RECTANGLE, ph_left, ph_top, ph_w, ph_h)
                ph.fill.background()
                ph.line.color.rgb = accent
                ph.line.width = Pt(1.5)
                tbox = sld.shapes.add_textbox(ph_left, ph_top + ph_h * 0.4, ph_w, Inches(0.8))
                tfph = tbox.text_frame
                tfp0 = tfph.paragraphs[0]
                tfp0.text = fig.get("placeholder_text") or "Figure / chart placeholder\n(export PNG and set figure.image_path)"
                tfp0.font.size = Pt(11)
                tfp0.font.color.rgb = body_rgb
                tfp0.alignment = PP_ALIGN.CENTER
                cap = fig.get("caption")
                if cap:
                    cb = sld.shapes.add_textbox(ph_left, ph_top + ph_h + Inches(0.08), ph_w, Inches(0.45))
                    ctf = cb.text_frame
                    cp = ctf.paragraphs[0]
                    cp.text = str(cap)
                    cp.font.size = Pt(10)
                    cp.font.color.rgb = body_rgb
                    cp.alignment = PP_ALIGN.CENTER
            else:
                _render_content_body(sld, left_margin, top_body, full_text_w, slide_h - Inches(1.65), slide,
                                     font_name=font_body, body_rgb=body_rgb, title_rgb=title_rgb, accent_rgb=accent)

        # ── Slide footer: separator line + page number (all non-title slides) ──
        if stype != "title":
            _sl_num = slide.get("slide_number", "")
            if _sl_num:
                _sep = sld.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.55), slide_h - Inches(0.26),
                    slide_w - Inches(1.1), Inches(0.007),
                )
                _sep.fill.solid()
                _sep.fill.fore_color.rgb = RGBColor(210, 215, 225)
                _sep.line.fill.background()
                _fn = sld.shapes.add_textbox(
                    slide_w - Inches(0.75), slide_h - Inches(0.24),
                    Inches(0.55), Inches(0.20),
                )
                _fnp = _fn.text_frame.paragraphs[0]
                _fnp.text = str(_sl_num)
                _fnp.font.size = Pt(9)
                _fnp.font.color.rgb = RGBColor(152, 160, 178)
                _fnp.alignment = PP_ALIGN.RIGHT
                try:
                    _fnp.font.name = font_body
                except Exception:
                    pass

        notes = slide.get("notes")
        if notes:
            try:
                ns = sld.notes_slide
                # Truncate notes to ~300 chars to stay within the notes pane
                notes_text = str(notes)
                if len(notes_text) > 300:
                    notes_text = notes_text[:297] + "..."
                ns.notes_text_frame.text = notes_text
            except Exception:
                pass

    Path(output_file).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_file)
    return f"Saved {len(slides_in)} slides to {output_file}"


def main() -> int:
    parser = argparse.ArgumentParser(description="Build editable summary PPTX from JSON plan")
    parser.add_argument("--plan-file", required=True, help="Path to presentation plan JSON")
    parser.add_argument("--output-file", required=True, help="Output .pptx path")
    args = parser.parse_args()

    with open(args.plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    msg = build_pptx(plan, args.output_file)
    print(msg)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
